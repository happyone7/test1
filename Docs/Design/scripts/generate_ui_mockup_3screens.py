#!/usr/bin/env python3
"""
UI_Mockup_3Screens.pptx 생성 스크립트
3화면: 타이틀 / 허브(스킬트리) / 인게임(HUD)
UI 에셋 16개 PNG 배치 상태 반영
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 프로젝트 루트
PROJECT_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..'))
ART_UI = os.path.join(PROJECT_ROOT, 'Assets', 'Art', 'UI')

# 색상 팔레트 (다크 판타지)
C_DEEP_BG    = RGBColor(0x0A, 0x0A, 0x12)
C_MAIN_BG    = RGBColor(0x14, 0x14, 0x20)
C_BRIGHT_BG  = RGBColor(0x1E, 0x1E, 0x30)
C_PANEL      = RGBColor(0x1A, 0x18, 0x28)
C_BRIGHT_PNL = RGBColor(0x24, 0x22, 0x36)
C_BORDER     = RGBColor(0x5A, 0x50, 0x70)
C_GOLD       = RGBColor(0xFF, 0xD8, 0x4D)
C_EMERALD    = RGBColor(0x40, 0xD4, 0x70)
C_SAPPHIRE   = RGBColor(0x40, 0x80, 0xD4)
C_RUBY       = RGBColor(0xD4, 0x40, 0x40)
C_PURPLE     = RGBColor(0x90, 0x60, 0xD0)
C_TEXT_MAIN  = RGBColor(0xE0, 0xDC, 0xD0)
C_TEXT_SUB   = RGBColor(0xA0, 0x98, 0x90)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# 슬라이드 크기 (16:9, 1920x1080에 해당)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def img_path(*parts):
    return os.path.join(ART_UI, *parts)


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    """사각형 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.color.rgb = border_color or C_BORDER
    shape.line.width = border_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=12,
             color=C_TEXT_MAIN, bold=False, alignment=PP_ALIGN.LEFT):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_image_safe(slide, path, left, top, width=None, height=None):
    """이미지 추가 (파일 없으면 건너뜀)"""
    if not os.path.exists(path):
        print(f"  [SKIP] {path}")
        return None
    kwargs = {'left': left, 'top': top}
    if width:
        kwargs['width'] = width
    if height:
        kwargs['height'] = height
    return slide.shapes.add_picture(path, **kwargs)


def create_title_slide(prs):
    """슬라이드 1: 타이틀 화면"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # 배경색
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_DEEP_BG

    # 배경 이미지 (전체)
    add_image_safe(slide, img_path('Backgrounds', 'TitleBG_01.png'),
                   Inches(0), Inches(0), width=SLIDE_W, height=SLIDE_H)

    # 로고 (중앙 상단)
    add_image_safe(slide, img_path('Logo', 'SoulspireLogo_02.png'),
                   Inches(4.2), Inches(0.8), width=Inches(5))

    # 버튼 영역 (중앙 하단)
    btn_x = Inches(4.5)
    btn_w = Inches(4.3)
    btn_h = Inches(0.7)
    btn_gap = Inches(0.15)

    # 시작 버튼 (accent)
    y = Inches(4.0)
    add_image_safe(slide, img_path('Buttons', 'btn_accent_idle.png'),
                   btn_x, y, width=btn_w, height=btn_h)
    add_text(slide, btn_x, y, btn_w, btn_h, "시작",
             font_size=20, color=C_EMERALD, bold=True, alignment=PP_ALIGN.CENTER)

    # 설정 버튼 (basic)
    y += btn_h + btn_gap
    add_image_safe(slide, img_path('Buttons', 'btn_basic_idle.png'),
                   btn_x, y, width=btn_w, height=btn_h)
    add_text(slide, btn_x, y, btn_w, btn_h, "설정",
             font_size=18, color=C_TEXT_MAIN, bold=False, alignment=PP_ALIGN.CENTER)

    # 종료 버튼 (basic)
    y += btn_h + btn_gap
    add_image_safe(slide, img_path('Buttons', 'btn_basic_idle.png'),
                   btn_x, y, width=btn_w, height=btn_h)
    add_text(slide, btn_x, y, btn_w, btn_h, "종료",
             font_size=18, color=C_RUBY, bold=False, alignment=PP_ALIGN.CENTER)

    # 에셋 사용 목록 주석
    add_text(slide, Inches(0.3), Inches(6.8), Inches(12), Inches(0.5),
             "[적용 에셋] TitleBG_01, SoulspireLogo_02, btn_accent_idle, btn_basic_idle x2",
             font_size=9, color=C_TEXT_SUB)


def create_hub_slide(prs):
    """슬라이드 2: 허브(스킬트리) 화면"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_DEEP_BG

    # 배경 이미지 (전체)
    add_image_safe(slide, img_path('Backgrounds', 'HubBG_03_dimmed.png'),
                   Inches(0), Inches(0), width=SLIDE_W, height=SLIDE_H)

    # === 상단 바 (50px -> ~0.52in) ===
    bar_h = Inches(0.52)
    add_image_safe(slide, img_path('Frames', 'panel_frame.png'),
                   Inches(0), Inches(0), width=SLIDE_W, height=bar_h)

    # Bit 아이콘 + 텍스트
    add_image_safe(slide, img_path('Icons', 'icon_bit.png'),
                   Inches(0.3), Inches(0.08), width=Inches(0.36), height=Inches(0.36))
    add_text(slide, Inches(0.7), Inches(0.05), Inches(1.5), bar_h,
             "Bit: 12,500", font_size=14, color=C_EMERALD, bold=True)

    # Core 아이콘 + 텍스트
    add_image_safe(slide, img_path('Icons', 'icon_core.png'),
                   Inches(2.5), Inches(0.08), width=Inches(0.36), height=Inches(0.36))
    add_text(slide, Inches(2.9), Inches(0.05), Inches(1.5), bar_h,
             "Core: 3", font_size=14, color=C_PURPLE, bold=True)

    # === 스킬 트리 영역 (중앙) ===
    tree_top = Inches(0.7)
    tree_h = Inches(5.5)

    # 스킬 노드 3개 (원형 배치)
    node_size = Inches(1.0)
    nodes = [
        (Inches(3.0), Inches(2.0), "공격력", "Lv.2/5", C_RUBY),
        (Inches(6.0), Inches(1.5), "공격속도", "Lv.1/5", C_SAPPHIRE),
        (Inches(9.0), Inches(2.5), "기지 HP", "Lv.3/5", C_EMERALD),
    ]
    for nx, ny, name, level, color in nodes:
        # 노드 배경 (tower_slot 활용)
        add_image_safe(slide, img_path('Frames', 'tower_slot.png'),
                       nx, ny, width=node_size, height=node_size)
        add_text(slide, nx, ny + Inches(0.25), node_size, Inches(0.3),
                 name, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, nx, ny + Inches(0.55), node_size, Inches(0.25),
                 level, font_size=9, color=C_TEXT_SUB, alignment=PP_ALIGN.CENTER)

    # === 상세 패널 (오른쪽) ===
    detail_x = Inches(8.5)
    detail_y = Inches(1.0)
    detail_w = Inches(4.5)
    detail_h = Inches(4.5)
    add_image_safe(slide, img_path('Frames', 'panel_frame.png'),
                   detail_x, detail_y, width=detail_w, height=detail_h)

    add_text(slide, detail_x + Inches(0.3), detail_y + Inches(0.2), Inches(3.5), Inches(0.4),
             "공격력 Lv2 -> Lv3", font_size=16, color=C_GOLD, bold=True)
    add_text(slide, detail_x + Inches(0.3), detail_y + Inches(0.7), Inches(3.5), Inches(0.3),
             "타워의 기본 공격력을 증가시킵니다.", font_size=11, color=C_TEXT_SUB)
    add_text(slide, detail_x + Inches(0.3), detail_y + Inches(1.2), Inches(3.5), Inches(0.3),
             "현재: +2.0   ->   변경 후: +3.0", font_size=12, color=C_TEXT_MAIN)
    add_text(slide, detail_x + Inches(0.3), detail_y + Inches(1.7), Inches(3.5), Inches(0.3),
             "비용: 500 Bit", font_size=14, color=C_GOLD, bold=True)

    # 구매 버튼
    add_image_safe(slide, img_path('Buttons', 'btn_accent_idle.png'),
                   detail_x + Inches(0.8), detail_y + Inches(2.3),
                   width=Inches(2.8), height=Inches(0.55))
    add_text(slide, detail_x + Inches(0.8), detail_y + Inches(2.3),
             Inches(2.8), Inches(0.55), "구매",
             font_size=16, color=C_EMERALD, bold=True, alignment=PP_ALIGN.CENTER)

    # === 하단 바 (55px -> ~0.57in) ===
    bot_y = Inches(6.93)
    bot_h = Inches(0.57)
    add_image_safe(slide, img_path('Frames', 'panel_frame.png'),
                   Inches(0), bot_y, width=SLIDE_W, height=bot_h)

    # 드롭다운 프레임
    add_image_safe(slide, img_path('Frames', 'dropdown_frame.png'),
                   Inches(0.5), bot_y + Inches(0.08), width=Inches(2.0), height=Inches(0.42))
    add_text(slide, Inches(0.7), bot_y + Inches(0.05), Inches(1.6), Inches(0.42),
             "Stage 1", font_size=12, color=C_TEXT_MAIN)

    # 출격 버튼
    add_image_safe(slide, img_path('Buttons', 'btn_accent_idle.png'),
                   Inches(4.5), bot_y + Inches(0.05), width=Inches(4.0), height=Inches(0.47))
    add_text(slide, Inches(4.5), bot_y + Inches(0.05), Inches(4.0), Inches(0.47),
             "출격!", font_size=16, color=C_EMERALD, bold=True, alignment=PP_ALIGN.CENTER)

    # 설정/종료
    add_image_safe(slide, img_path('Buttons', 'btn_basic_idle.png'),
                   Inches(9.5), bot_y + Inches(0.05), width=Inches(1.8), height=Inches(0.47))
    add_text(slide, Inches(9.5), bot_y + Inches(0.05), Inches(1.8), Inches(0.47),
             "설정", font_size=12, color=C_TEXT_MAIN, alignment=PP_ALIGN.CENTER)

    add_image_safe(slide, img_path('Buttons', 'btn_basic_idle.png'),
                   Inches(11.5), bot_y + Inches(0.05), width=Inches(1.5), height=Inches(0.47))
    add_text(slide, Inches(11.5), bot_y + Inches(0.05), Inches(1.5), Inches(0.47),
             "종료", font_size=12, color=C_RUBY, alignment=PP_ALIGN.CENTER)

    # 에셋 목록 주석
    add_text(slide, Inches(0.3), Inches(6.3), Inches(12), Inches(0.4),
             "[적용 에셋] HubBG_03_dimmed, panel_frame x3, icon_bit, icon_core, tower_slot x3, dropdown_frame, btn_accent_idle x2, btn_basic_idle x2",
             font_size=8, color=C_TEXT_SUB)


def create_ingame_slide(prs):
    """슬라이드 3: 인게임(HUD) 화면"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_MAIN_BG

    # === 상단 HUD 바 ===
    hud_h = Inches(0.65)
    add_image_safe(slide, img_path('Frames', 'panel_frame.png'),
                   Inches(0), Inches(0), width=SLIDE_W, height=hud_h)

    # 웨이브 텍스트
    add_text(slide, Inches(0.3), Inches(0.1), Inches(2), Inches(0.45),
             "웨이브: 3/5", font_size=14, color=C_TEXT_MAIN, bold=True)

    # Bit 카운터
    add_image_safe(slide, img_path('Icons', 'icon_bit.png'),
                   Inches(3.0), Inches(0.1), width=Inches(0.36), height=Inches(0.36))
    add_text(slide, Inches(3.4), Inches(0.1), Inches(1.5), Inches(0.45),
             "Bit: 2,340", font_size=13, color=C_EMERALD, bold=True)

    # 기지 HP 바
    add_text(slide, Inches(6.0), Inches(0.05), Inches(1.5), Inches(0.3),
             "기지 HP:", font_size=11, color=C_TEXT_SUB)

    # HP 바 프레임
    hp_x = Inches(7.5)
    hp_w = Inches(4.5)
    hp_h = Inches(0.35)
    add_image_safe(slide, img_path('Frames', 'hp_bar_frame.png'),
                   hp_x, Inches(0.12), width=hp_w, height=hp_h)

    # HP 바 채우기 (80%)
    fill_w = Inches(4.5 * 0.8)
    add_image_safe(slide, img_path('Frames', 'hp_bar_fill.png'),
                   hp_x + Inches(0.03), Inches(0.15), width=fill_w, height=Inches(0.28))

    add_text(slide, hp_x, Inches(0.12), hp_w, hp_h,
             "16/20", font_size=11, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # === 중앙 게임 영역 (빈 공간, 설명 텍스트) ===
    add_text(slide, Inches(4.5), Inches(3.2), Inches(4.5), Inches(0.5),
             "[ 게임 플레이 영역 ]", font_size=16, color=C_BORDER, alignment=PP_ALIGN.CENTER)

    # === 배속/일시정지 버튼 (우측 하단) ===
    speed_y = Inches(5.8)
    speed_w = Inches(0.8)
    speed_h = Inches(0.45)
    for i, label in enumerate(["x1", "x2", "x3"]):
        x = Inches(9.5) + Inches(i * 1.0)
        add_image_safe(slide, img_path('Buttons', 'btn_basic_idle.png'),
                       x, speed_y, width=speed_w, height=speed_h)
        color = C_EMERALD if i == 0 else C_TEXT_SUB
        add_text(slide, x, speed_y, speed_w, speed_h,
                 label, font_size=11, color=color, bold=(i == 0), alignment=PP_ALIGN.CENTER)

    # 일시정지
    add_image_safe(slide, img_path('Buttons', 'btn_basic_idle.png'),
                   Inches(12.5), speed_y, width=speed_w, height=speed_h)
    add_text(slide, Inches(12.5), speed_y, speed_w, speed_h,
             "||", font_size=12, color=C_TEXT_MAIN, alignment=PP_ALIGN.CENTER)

    # === 하단 타워 인벤토리 바 ===
    inv_y = Inches(6.5)
    inv_h = Inches(0.85)
    inv_w = Inches(8.0)
    inv_x = Inches(2.7)

    # 바 배경 (panel_frame)
    add_image_safe(slide, img_path('Frames', 'panel_frame.png'),
                   inv_x, inv_y, width=inv_w, height=inv_h)

    # 8개 슬롯 (tower_slot)
    slot_size = Inches(0.7)
    slot_gap = Inches(0.12)
    slots_start_x = inv_x + Inches(0.35)
    for i in range(8):
        sx = slots_start_x + (slot_size + slot_gap) * i
        add_image_safe(slide, img_path('Frames', 'tower_slot.png'),
                       sx, inv_y + Inches(0.08), width=slot_size, height=slot_size)
        if i < 3:
            # 채워진 슬롯 예시
            add_text(slide, sx, inv_y + Inches(0.52), slot_size, Inches(0.2),
                     f"Lv.{i + 1}", font_size=8, color=C_TEXT_SUB, alignment=PP_ALIGN.CENTER)

    # === 타워 정보 툴팁 (우측에 표시 예시) ===
    tt_x = Inches(9.0)
    tt_y = Inches(3.0)
    tt_w = Inches(3.5)
    tt_h = Inches(2.5)
    add_image_safe(slide, img_path('Frames', 'tooltip_frame.png'),
                   tt_x, tt_y, width=tt_w, height=tt_h)

    add_text(slide, tt_x + Inches(0.2), tt_y + Inches(0.15), tt_w, Inches(0.3),
             "화살탑  Lv.2", font_size=14, color=C_SAPPHIRE, bold=True)
    add_text(slide, tt_x + Inches(0.2), tt_y + Inches(0.55), Inches(2.5), Inches(0.25),
             "공격력:  12.5", font_size=11, color=C_TEXT_MAIN)
    add_text(slide, tt_x + Inches(0.2), tt_y + Inches(0.85), Inches(2.5), Inches(0.25),
             "공격속도:  1.20/s", font_size=11, color=C_TEXT_MAIN)
    add_text(slide, tt_x + Inches(0.2), tt_y + Inches(1.15), Inches(2.5), Inches(0.25),
             "사거리:  4.0", font_size=11, color=C_TEXT_MAIN)
    add_text(slide, tt_x + Inches(0.2), tt_y + Inches(1.55), Inches(2.5), Inches(0.25),
             "판매: +25 Bit", font_size=10, color=RGBColor(0xE0, 0x80, 0x30))

    # 에셋 목록 주석
    add_text(slide, Inches(0.3), Inches(7.1), Inches(12), Inches(0.4),
             "[적용 에셋] panel_frame x2, hp_bar_frame, hp_bar_fill, icon_bit, tower_slot x8, tooltip_frame, btn_basic_idle x5",
             font_size=8, color=C_TEXT_SUB)


def create_asset_summary_slide(prs):
    """슬라이드 4: 에셋 매핑 요약"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_DEEP_BG

    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "UI 에셋 매핑 요약 (16개 + 배경/로고 3개)", font_size=24, color=C_GOLD, bold=True)

    mapping = [
        ("btn_basic_idle/hover/pressed/disabled", "타이틀(설정/종료), 허브(설정/종료), 인게임(배속/일시정지/Hub)"),
        ("btn_accent_idle/hover/pressed/disabled", "타이틀(시작), 허브(출격/구매), 인게임(재도전)"),
        ("panel_frame", "허브(상단바/하단바/상세패널), 인게임(HUD바/인벤토리바/런종료패널), 설정팝업, 방치팝업"),
        ("hp_bar_frame", "인게임 기지 HP바 배경"),
        ("hp_bar_fill", "인게임 기지 HP바 채우기"),
        ("tower_slot", "허브(스킬노드), 인게임(인벤토리 슬롯 x8)"),
        ("tooltip_frame", "인게임 타워 정보 툴팁"),
        ("dropdown_frame", "허브 스테이지 선택 드롭다운"),
        ("icon_bit", "허브(상단바), 인게임(HUD)"),
        ("icon_core", "허브(상단바)"),
        ("TitleBG_01", "타이틀 화면 배경"),
        ("HubBG_03_dimmed", "허브 화면 배경"),
        ("SoulspireLogo_02", "타이틀 화면 로고"),
    ]

    y = Inches(1.2)
    for asset, usage in mapping:
        add_text(slide, Inches(0.5), y, Inches(3.5), Inches(0.3),
                 asset, font_size=10, color=C_EMERALD, bold=True)
        add_text(slide, Inches(4.2), y, Inches(8.5), Inches(0.3),
                 usage, font_size=10, color=C_TEXT_MAIN)
        y += Inches(0.38)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("=== UI Mockup 3Screens PPT 생성 ===")
    print("슬라이드 1: 타이틀 화면...")
    create_title_slide(prs)

    print("슬라이드 2: 허브(스킬트리) 화면...")
    create_hub_slide(prs)

    print("슬라이드 3: 인게임(HUD) 화면...")
    create_ingame_slide(prs)

    print("슬라이드 4: 에셋 매핑 요약...")
    create_asset_summary_slide(prs)

    output = os.path.join(PROJECT_ROOT, 'Docs', 'Design', 'UI_Mockup_3Screens.pptx')
    prs.save(output)
    print(f"\n[완료] {output}")


if __name__ == '__main__':
    main()
