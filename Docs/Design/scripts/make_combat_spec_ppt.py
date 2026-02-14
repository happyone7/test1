#!/usr/bin/env python3
"""
Soulspire - 인게임 전투 시스템 기획서 PPT 생성기
다크 판타지 테마 / 해상도: 1920x1080 (16:9)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── 상수 ─────────────────────────────────────────────────────────────────
OUT_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "InGameCombat_v0.2.pptx",
)

# 1 px = 9525 EMU (96 DPI 기준)
PX = 9525

# 슬라이드 크기 (1920x1080 px)
SLIDE_W = 1920 * PX
SLIDE_H = 1080 * PX

# ─── 다크 판타지 컬러 팔레트 ─────────────────────────────────────────────
# 배경
C_BG_DEEP     = RGBColor(0x0A, 0x0A, 0x12)   # #0A0A12 최심부
C_BG_MAIN     = RGBColor(0x14, 0x14, 0x20)   # #141420 기본 배경
C_BG_LIGHT    = RGBColor(0x1E, 0x1E, 0x30)   # #1E1E30 밝은 배경

# UI
C_UI_BG       = RGBColor(0x12, 0x10, 0x1A)   # #12101A 패널 배경
C_UI_PANEL    = RGBColor(0x1A, 0x18, 0x28)   # #1A1828 기본 패널
C_UI_PANEL_LT = RGBColor(0x24, 0x22, 0x36)   # #242236 밝은 패널

# 프레임/테두리
C_FRAME       = RGBColor(0x5A, 0x50, 0x70)   # #5A5070 기본 프레임
C_GOLD        = RGBColor(0xB0, 0xA0, 0x80)   # #B0A080 금박
C_GOLD_BRIGHT = RGBColor(0xFF, 0xD8, 0x4D)   # #FFD84D 순금 강조

# 텍스트
C_TEXT_MAIN   = RGBColor(0xE0, 0xDC, 0xD0)   # #E0DCD0 메인
C_TEXT_SUB    = RGBColor(0xA0, 0x98, 0x90)   # #A09890 서브
C_TITLE       = RGBColor(0xE8, 0xE4, 0xF0)   # #E8E4F0 제목 (백색 마법)

# 마법 하이라이트
C_RUBY        = RGBColor(0xD4, 0x40, 0x40)   # #D44040
C_EMERALD     = RGBColor(0x40, 0xD4, 0x70)   # #40D470
C_SAPPHIRE    = RGBColor(0x40, 0x80, 0xD4)   # #4080D4
C_PURPLE      = RGBColor(0x90, 0x60, 0xD0)   # #9060D0
C_FLAME       = RGBColor(0xE0, 0x80, 0x30)   # #E08030
C_LIGHTNING   = RGBColor(0xE0, 0xD0, 0x40)   # #E0D040

# 타워 시그니처 컬러
C_TOWER_ARROW     = RGBColor(0xC0, 0xA8, 0x70)   # #C0A870
C_TOWER_CANNON    = RGBColor(0xE0, 0x80, 0x30)   # #E08030
C_TOWER_ICE       = RGBColor(0x40, 0x80, 0xD4)   # #4080D4
C_TOWER_LIGHTNING = RGBColor(0xE0, 0xD0, 0x40)   # #E0D040
C_TOWER_LASER     = RGBColor(0x90, 0x60, 0xD0)   # #9060D0
C_TOWER_VOID      = RGBColor(0x30, 0x20, 0x50)   # #302050

C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def px(val):
    """픽셀 값을 EMU로 변환."""
    return int(val * PX)


def add_bg(slide, color=C_BG_DEEP):
    """슬라이드 배경을 단색으로 채움."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=None, border_color=None, border_width=1):
    """사각형 도형 추가."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, x, y, w, h, fill_color=None, border_color=None, border_width=1):
    """둥근 사각형 도형 추가."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, font_size=14, color=C_TEXT_MAIN,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """텍스트 상자 추가."""
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    tf.auto_size = None
    return txBox


def add_multiline_text(slide, x, y, w, h, lines, font_size=12, color=C_TEXT_MAIN,
                       bold=False, align=PP_ALIGN.LEFT):
    """여러 줄 텍스트 상자 추가. lines는 (텍스트, 크기, 색상, 굵기) 튜플 또는 str."""
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            txt, sz, clr, b = line_data, font_size, color, bold
        else:
            txt = line_data[0]
            sz = line_data[1] if len(line_data) > 1 else font_size
            clr = line_data[2] if len(line_data) > 2 else color
            b = line_data[3] if len(line_data) > 3 else bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = b
        p.alignment = align
    return txBox


def add_arrow_shape(slide, x1, y1, x2, y2, color=C_GOLD, width=2):
    """두 점 사이에 화살표 커넥터를 그린다 (얇은 사각형 + 삼각형 화살표 머리)."""
    # 간단하게: 가로 화살표만 지원 (y1 == y2)
    if y1 == y2:
        # 수평 선
        line_y = y1
        line_x = min(x1, x2)
        line_w = abs(x2 - x1)
        add_rect(slide, line_x, line_y, line_w, 3, fill_color=color)
        # 화살표 머리 (작은 삼각형 근사 - 사각형)
        if x2 > x1:
            add_rect(slide, x2 - 6, line_y - 4, 6, 11, fill_color=color)
        else:
            add_rect(slide, x1 - 6, line_y - 4, 6, 11, fill_color=color)
    elif x1 == x2:
        # 수직 선
        line_x = x1
        line_y = min(y1, y2)
        line_h = abs(y2 - y1)
        add_rect(slide, line_x, line_y, 3, line_h, fill_color=color)
        if y2 > y1:
            add_rect(slide, line_x - 4, y2 - 6, 11, 6, fill_color=color)
        else:
            add_rect(slide, line_x - 4, y1 - 6, 11, 6, fill_color=color)


def add_gold_border(slide, x, y, w, h, border_width=2):
    """금박 테두리를 추가한다."""
    add_rect(slide, x, y, w, h, fill_color=None, border_color=C_GOLD, border_width=border_width)


def add_panel(slide, x, y, w, h, title=None, title_size=20):
    """다크 판타지 스타일 패널 (배경 + 프레임 테두리)."""
    add_rounded_rect(slide, x, y, w, h,
                     fill_color=C_UI_PANEL, border_color=C_FRAME, border_width=2)
    if title:
        add_text(slide, x + 15, y + 10, w - 30, 30, title,
                 font_size=title_size, color=C_TITLE, bold=True)
        # 구분선
        add_rect(slide, x + 15, y + 40, w - 30, 1, fill_color=C_FRAME)
    return


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DEEP)

    # 금박 외곽 테두리
    add_rect(slide, 30, 30, 1860, 1020,
             fill_color=None, border_color=C_GOLD, border_width=3)
    # 내부 얇은 테두리
    add_rect(slide, 45, 45, 1830, 990,
             fill_color=None, border_color=C_FRAME, border_width=1)

    # 장식 상단 라인
    add_rect(slide, 300, 250, 1320, 2, fill_color=C_GOLD)

    # 게임 타이틀
    add_text(slide, 260, 280, 1400, 100,
             "SOULSPIRE", font_size=56, color=C_GOLD_BRIGHT,
             bold=True, align=PP_ALIGN.CENTER)

    # 부제목
    add_text(slide, 260, 390, 1400, 60,
             "인게임 전투 시스템 기획서 v0.2",
             font_size=24, color=C_TITLE, bold=True, align=PP_ALIGN.CENTER)

    # 장식 하단 라인
    add_rect(slide, 300, 470, 1320, 2, fill_color=C_GOLD)

    # 메타 정보
    meta_lines = [
        ("", 10, C_TEXT_SUB, False),
        ("장르: 다크 판타지 타워 디펜스 (로그라이트)", 14, C_TEXT_SUB, False),
        ("대상 독자: 프로그래밍팀장, QA팀장", 14, C_TEXT_SUB, False),
        ("최종 수정: 2026-02-15", 14, C_TEXT_SUB, False),
        ("작성자: 기획팀장 (Game Designer)", 14, C_TEXT_SUB, False),
        ("상태: 총괄PD 피드백 반영 완료", 14, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 560, 510, 800, 220, meta_lines, align=PP_ALIGN.CENTER)

    # 하단 장식
    add_rect(slide, 600, 800, 720, 1, fill_color=C_FRAME)
    add_text(slide, 560, 810, 800, 30,
             "Dark Fantasy  ·  Pixel Art  ·  3/4 Quarter View",
             font_size=11, color=C_FRAME, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 2: 런 플로우 다이어그램
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_run_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DEEP)

    # 제목
    add_text(slide, 60, 30, 800, 40, "1. 런 플로우 (상태 전환)",
             font_size=22, color=C_TITLE, bold=True)
    add_rect(slide, 60, 70, 400, 2, fill_color=C_GOLD)

    # ── 메인 플로우 (가로) ──
    # 박스 좌표
    box_w, box_h = 200, 70
    y_main = 200

    # Hub
    add_rounded_rect(slide, 80, y_main, box_w, box_h,
                     fill_color=C_UI_PANEL, border_color=C_SAPPHIRE, border_width=2)
    add_text(slide, 80, y_main + 10, box_w, 25, "Hub (성소)",
             font_size=16, color=C_SAPPHIRE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 80, y_main + 38, box_w, 25, "스킬 트리 / 보물상자",
             font_size=10, color=C_TEXT_SUB, align=PP_ALIGN.CENTER)

    # 화살표 Hub -> InGame:준비
    add_rect(slide, 280, y_main + 33, 80, 3, fill_color=C_GOLD)
    add_text(slide, 285, y_main - 5, 70, 20, "출격",
             font_size=9, color=C_GOLD, align=PP_ALIGN.CENTER)

    # InGame: 준비
    add_rounded_rect(slide, 360, y_main, box_w, box_h,
                     fill_color=C_UI_PANEL, border_color=C_EMERALD, border_width=2)
    add_text(slide, 360, y_main + 10, box_w, 25, "InGame: 준비",
             font_size=16, color=C_EMERALD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 360, y_main + 38, box_w, 25, "기지HP초기화/타워복원",
             font_size=10, color=C_TEXT_SUB, align=PP_ALIGN.CENTER)

    # 화살표
    add_rect(slide, 560, y_main + 33, 80, 3, fill_color=C_GOLD)

    # InGame: 전투
    add_rounded_rect(slide, 640, y_main, 240, box_h,
                     fill_color=C_UI_PANEL, border_color=C_FLAME, border_width=2)
    add_text(slide, 640, y_main + 10, 240, 25, "InGame: 전투",
             font_size=16, color=C_FLAME, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 640, y_main + 38, 240, 25, "웨이브 스폰 / 타워 공격",
             font_size=10, color=C_TEXT_SUB, align=PP_ALIGN.CENTER)

    # 분기 화살표 (위: 클리어, 아래: 패배)
    # 승리 경로 (위로)
    y_win = y_main - 100
    add_rect(slide, 880, y_main + 20, 120, 3, fill_color=C_EMERALD)
    add_text(slide, 885, y_main + 2, 110, 18, "최종웨이브 클리어",
             font_size=8, color=C_EMERALD, align=PP_ALIGN.CENTER)

    # RunEnd: 승리
    add_rounded_rect(slide, 1000, y_main - 20, box_w + 40, box_h,
                     fill_color=C_UI_PANEL, border_color=C_EMERALD, border_width=2)
    add_text(slide, 1000, y_main - 10, box_w + 40, 25, "RunEnd: 승리",
             font_size=16, color=C_EMERALD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 1000, y_main + 18, box_w + 40, 25, "Bit+Core 정산/다음스테이지",
             font_size=10, color=C_TEXT_SUB, align=PP_ALIGN.CENTER)

    # 패배 경로 (아래로)
    y_lose = y_main + 130
    add_rect(slide, 760, y_main + 70, 3, 60, fill_color=C_RUBY)
    add_text(slide, 770, y_main + 80, 100, 18, "기지 HP 0",
             font_size=8, color=C_RUBY)
    add_rect(slide, 760, y_lose + 33, 240, 3, fill_color=C_RUBY)

    # RunEnd: 패배
    add_rounded_rect(slide, 1000, y_lose, box_w + 40, box_h,
                     fill_color=C_UI_PANEL, border_color=C_RUBY, border_width=2)
    add_text(slide, 1000, y_lose + 10, box_w + 40, 25, "RunEnd: 패배",
             font_size=16, color=C_RUBY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 1000, y_lose + 38, box_w + 40, 25, "Bit 정산 / 재도전",
             font_size=10, color=C_TEXT_SUB, align=PP_ALIGN.CENTER)

    # 복귀 화살표 (RunEnd -> Hub)
    add_rect(slide, 1240, y_main + 15, 120, 3, fill_color=C_GOLD)
    add_rect(slide, 1240, y_lose + 35, 120, 3, fill_color=C_GOLD)
    # Hub 복귀 박스 (우측)
    add_rounded_rect(slide, 1360, y_main + 40, 160, 50,
                     fill_color=C_UI_PANEL, border_color=C_SAPPHIRE, border_width=2)
    add_text(slide, 1360, y_main + 50, 160, 30, "Hub 복귀",
             font_size=14, color=C_SAPPHIRE, bold=True, align=PP_ALIGN.CENTER)

    # 수직 연결선
    add_rect(slide, 1440, y_main + 15, 3, 27, fill_color=C_GOLD)
    add_rect(slide, 1440, y_main + 90, 3, y_lose - y_main - 55, fill_color=C_GOLD)
    add_rect(slide, 1360, y_lose + 35, 80, 3, fill_color=C_GOLD)

    # ── 하단: 웨이브 반복 루프 ──
    loop_y = 480
    add_text(slide, 60, loop_y - 20, 600, 30, "웨이브 진행 루프",
             font_size=18, color=C_TITLE, bold=True)
    add_rect(slide, 60, loop_y + 15, 300, 2, fill_color=C_GOLD)

    loop_box_w, loop_box_h = 190, 55
    ly = loop_y + 50

    boxes = [
        ("웨이브 N 스폰", C_FLAME, 80),
        ("모든 Node 처리", C_EMERALD, 340),
        ("웨이브 클리어 판정", C_GOLD_BRIGHT, 600),
        ("다음 웨이브 대기", C_SAPPHIRE, 860),
    ]
    for label, clr, bx in boxes:
        add_rounded_rect(slide, bx, ly, loop_box_w, loop_box_h,
                         fill_color=C_UI_PANEL, border_color=clr, border_width=2)
        add_text(slide, bx, ly + 13, loop_box_w, 30, label,
                 font_size=13, color=clr, bold=True, align=PP_ALIGN.CENTER)

    # 화살표들
    for i in range(3):
        sx = boxes[i][2] + loop_box_w
        ex = boxes[i + 1][2]
        add_rect(slide, sx, ly + 26, ex - sx, 3, fill_color=C_GOLD)

    # 루프 리턴 화살표 (아래쪽으로 돌아감)
    add_rect(slide, boxes[3][2] + loop_box_w, ly + 26, 50, 3, fill_color=C_GOLD)
    add_rect(slide, boxes[3][2] + loop_box_w + 47, ly + 26, 3, 60, fill_color=C_GOLD)
    add_rect(slide, 80, ly + 83, boxes[3][2] + loop_box_w - 30, 3, fill_color=C_GOLD)
    add_rect(slide, 80, ly + 26, 3, 60, fill_color=C_GOLD)

    # ── FTUE 안내 ──
    ftue_y = 680
    add_panel(slide, 60, ftue_y, 700, 130, title="첫 런 특수 처리 (FTUE)")
    ftue_lines = [
        ("• 게임 최초 시작 시 Hub 스킵 → 바로 스테이지 1 진입", 12, C_TEXT_MAIN, False),
        ("• Arrow Tower 1기가 맵 정가운데 자동 배치", 12, C_TEXT_MAIN, False),
        ("• 웨이브 1 시작 전 3초 딜레이 (맵 확인 시간)", 12, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 80, ftue_y + 50, 660, 70, ftue_lines)

    # ── 런 종료 조건 ──
    add_panel(slide, 800, ftue_y, 700, 130, title="런 종료 조건")
    end_lines = [
        ("패배: 기지 HP 0 → 남은 Node 페이드아웃 → 보상 정산", 12, C_RUBY, False),
        ("승리: 최종 웨이브 클리어 → 클리어 연출 1.5초 → Core 획득", 12, C_EMERALD, False),
        ("공통: 타워 배치 SaveData 저장 (크래시 대비 즉시 저장)", 12, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 820, ftue_y + 50, 660, 70, end_lines)

    # ── 런 종료 패널 버튼 ──
    btn_y = ftue_y + 160
    add_panel(slide, 60, btn_y, 1440, 100, title="RunEnd 패널 버튼 동작")
    btn_lines = [
        ("패배 → Hub 버튼: GoToHub() | 재도전 버튼: StartRun(같은 스테이지)", 12, C_TEXT_MAIN, False),
        ("승리 → Hub 버튼: GoToHub() | 다음 스테이지 버튼: StartRun(다음 스테이지)", 12, C_TEXT_MAIN, False),
    ]
    add_multiline_text(slide, 80, btn_y + 48, 1400, 50, btn_lines)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 3: 타워 배치/합성
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_tower_placement(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DEEP)

    add_text(slide, 60, 30, 800, 40, "2. 타워 배치 / 합성 시스템",
             font_size=22, color=C_TITLE, bold=True)
    add_rect(slide, 60, 70, 400, 2, fill_color=C_GOLD)

    # ── 격자 배치 시각화 ──
    add_panel(slide, 60, 100, 580, 400, title="격자 기반 배치")

    grid_x0, grid_y0 = 90, 160
    cell = 45
    cols, rows = 8, 6
    labels_map = {}
    # W = 벽, P = 경로, T = 배치 가능
    grid_data = [
        ['W','W','W','W','W','W','W','W'],
        ['W','P','P','P','P','P','P','W'],
        ['W','T','T','P','T','T','P','W'],
        ['W','P','P','P','P','P','P','W'],
        ['W','T','P','T','T','P','T','W'],
        ['W','W','W','W','W','W','W','W'],
    ]
    cell_colors = {
        'W': (RGBColor(0x2A, 0x2A, 0x3A), C_FRAME),
        'P': (RGBColor(0x1E, 0x1E, 0x30), C_FRAME),
        'T': (RGBColor(0x15, 0x25, 0x18), C_EMERALD),
    }
    for r in range(rows):
        for c in range(cols):
            t = grid_data[r][c]
            fc, bc = cell_colors[t]
            add_rect(slide, grid_x0 + c * cell, grid_y0 + r * cell,
                     cell - 1, cell - 1, fill_color=fc, border_color=bc, border_width=1)
            if t == 'T':
                add_text(slide, grid_x0 + c * cell, grid_y0 + r * cell + 12,
                         cell - 1, 20, "T", font_size=8, color=C_EMERALD,
                         align=PP_ALIGN.CENTER)
            elif t == 'P':
                add_text(slide, grid_x0 + c * cell, grid_y0 + r * cell + 12,
                         cell - 1, 20, "P", font_size=8, color=C_TEXT_SUB,
                         align=PP_ALIGN.CENTER)

    # 범례
    legend_y = grid_y0 + rows * cell + 10
    legend_items = [
        ("W = 벽 (배치 불가)", C_FRAME),
        ("P = 경로 (Node 이동)", C_TEXT_SUB),
        ("T = 배치 가능 칸", C_EMERALD),
    ]
    for i, (txt, clr) in enumerate(legend_items):
        add_text(slide, 90, legend_y + i * 22, 300, 20, txt,
                 font_size=10, color=clr)

    # 배치 명세
    spec_y = legend_y + 75
    spec_lines = [
        ("타일 크기: 1x1 Unity Unit (32x32px, PPU 32)", 10, C_TEXT_SUB, False),
        ("초기 배치 칸: 12칸 (스테이지 1)", 10, C_TEXT_SUB, False),
        ("최대 확장: +4칸/회 (스킬 트리, 최대 3회)", 10, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 90, spec_y, 530, 60, spec_lines)

    # ── 드래그 배치 흐름 ──
    add_panel(slide, 660, 100, 700, 220, title="드래그 배치 UX")
    drag_lines = [
        ("1. 하단 인벤토리에서 타워 아이콘 터치/클릭", 12, C_TEXT_MAIN, False),
        ("2. 아이콘이 커서를 따라 이동 (드래그 상태)", 12, C_TEXT_MAIN, False),
        ("3. 맵 위 이동 시 하이라이트:", 12, C_TEXT_MAIN, False),
        ("   · 배치 가능 빈 칸 → 초록색 하이라이트", 11, C_EMERALD, False),
        ("   · 같은 타입+레벨 타워 위 → 금색 + 'LV UP!'", 11, C_GOLD_BRIGHT, False),
        ("   · 불가 위치 → 빨간색 하이라이트", 11, C_RUBY, False),
        ("4. 드롭 → 배치/합성 실행 또는 인벤토리 복귀", 12, C_TEXT_MAIN, False),
    ]
    add_multiline_text(slide, 680, 150, 660, 160, drag_lines)

    # ── 합성 규칙 ──
    add_panel(slide, 660, 340, 700, 160, title="합성 규칙")
    merge_lines = [
        ("조건: 같은 타입 + 같은 레벨 타워 위에 드롭", 12, C_TEXT_MAIN, False),
        ("결과: Level + 1 (예: Lv1+Lv1=Lv2, Lv2+Lv2=Lv3)", 12, C_TEXT_MAIN, False),
        ("최대 레벨: Lv4 (Lv4에는 합성 불가 → 빨간색 표시)", 12, C_GOLD_BRIGHT, True),
        ("합성 비용: 없음 (타워 1기 소모가 유일한 비용)", 12, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 680, 390, 660, 100, merge_lines)

    # ── 합성 이펙트 ──
    eff_y = 520
    add_panel(slide, 60, eff_y, 580, 130, title="합성 이펙트 스케일링")
    eff_lines = [
        ("Lv1→Lv2: 작은 빛남 + '+1'", 11, C_TEXT_SUB, False),
        ("Lv2→Lv3: 중간 빛남 + 짧은 파티클", 11, C_TEXT_MAIN, False),
        ("Lv3→Lv4: 화면 미세 흔들림 + 전용 발광 오라 + 'MAX!'", 11, C_GOLD_BRIGHT, True),
    ]
    add_multiline_text(slide, 80, eff_y + 48, 540, 70, eff_lines)

    # ── 핵심 규칙 패널 ──
    rule_y = 670
    add_panel(slide, 60, rule_y, 1300, 170, title="핵심 규칙")
    rule_lines = [
        ("타워 판매 불가 (총괄PD 확정)", 14, C_RUBY, True),
        ("  → 한번 배치하면 영구 점유. 배치 결정에 무게감 부여.", 12, C_TEXT_SUB, False),
        ("", 6, C_TEXT_SUB, False),
        ("배치 유지: 런이 종료되어도(죽어도) 배치된 타워는 유지", 14, C_EMERALD, True),
        ("  → 성장이 누적된다는 느낌의 핵심. SaveData에 즉시 저장.", 12, C_TEXT_SUB, False),
        ("", 6, C_TEXT_SUB, False),
        ("스테이지별 독립 배치: 각 스테이지는 독립적인 배치 레이아웃", 14, C_SAPPHIRE, True),
    ]
    add_multiline_text(slide, 80, rule_y + 48, 1260, 110, rule_lines)

    # ── 인벤토리 ──
    inv_y = 860
    add_panel(slide, 60, inv_y, 1300, 100, title="타워 인벤토리")
    inv_lines = [
        ("최대 8칸 | 타워 아이콘+이름+레벨 표시 | 타입별 그룹핑, 레벨순 정렬", 12, C_TEXT_MAIN, False),
        ("인벤토리 풀 → 보물상자 열기 불가 / 타워 획득: 오직 Hub 보물상자 + 첫 런 자동지급", 12, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 80, inv_y + 48, 1260, 40, inv_lines)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 4: 타워 6종 요약표
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_towers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DEEP)

    add_text(slide, 60, 30, 800, 40, "3. 타워 6종 요약",
             font_size=22, color=C_TITLE, bold=True)
    add_rect(slide, 60, 70, 400, 2, fill_color=C_GOLD)

    # 테이블 헤더
    hdr_y = 100
    col_x = [60, 220, 470, 650, 850, 1050, 1350]
    col_w = [160, 250, 180, 200, 200, 300, 200]
    headers = ["타워", "Attack 방식", "투사체", "타겟팅", "특수 데이터", "레벨별 핵심 스탯", "시그니처"]
    for i, h in enumerate(headers):
        add_rect(slide, col_x[i], hdr_y, col_w[i], 35,
                 fill_color=C_UI_PANEL, border_color=C_GOLD, border_width=1)
        add_text(slide, col_x[i] + 5, hdr_y + 5, col_w[i] - 10, 25, h,
                 font_size=11, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)

    towers = [
        ("Arrow\n마법 석궁탑", "투사체 단일 타격", "Projectile\n(기존)", "가장 가까운 적", "-",
         "DMG 10/18/30/50\nSPD 1.0/1.15/1.3/1.5", C_TOWER_ARROW),
        ("Cannon\n화염 마법포", "투사체 → 폭발 AoE", "CannonProjectile\n(신규)", "가장 가까운 적",
         "explosionRadius\n1.0/1.2/1.4/1.6",
         "DMG 20/40/70/110\n중심100%~가장자리50%", C_TOWER_CANNON),
        ("Ice\n빙결 마탑", "투사체 + 감속", "IceProjectile\n(신규)", "가장 가까운 적",
         "slow 25~65%\ndur 1.5~2.5s",
         "DMG 8/15/28/45\n감속 중첩불가,덮어쓰기", C_TOWER_ICE),
        ("Lightning\n뇌전 첨탑", "즉발 체인", "없음 (즉발)", "가장 가까운 적\n→ 체인",
         "chain 2~6\ndecay 30~10%",
         "DMG 20/32/48/70\n체인마다 데미지 감소", C_TOWER_LIGHTNING),
        ("Laser\n마력 포탑", "빔 지속 → 쿨다운", "없음\n(LineRenderer)", "가장 가까운 적",
         "beam 2~3s\nCD 3~2s, pierce",
         "DPS 25/50/90/140\n빔 관통 3~무제한", C_TOWER_LASER),
        ("Void\n공허의 오벨리스크", "범위 지속 데미지", "없음\n(상시 AoE)", "범위 내 전체",
         "AoE 2.0~3.5\nLv3+ 방어무시",
         "DPS 15/35/65/100\nLv4+ 감속 20%", C_TOWER_VOID),
    ]

    row_h = 130
    for i, (name, atk, proj, tgt, special, stats, clr) in enumerate(towers):
        ry = hdr_y + 35 + i * row_h
        # 타워 이름 셀 (시그니처 컬러 배경)
        add_rect(slide, col_x[0], ry, col_w[0], row_h,
                 fill_color=C_UI_BG, border_color=clr, border_width=2)
        add_multiline_text(slide, col_x[0] + 5, ry + 8, col_w[0] - 10, row_h - 10,
                           [(name.split('\n')[0], 13, clr, True),
                            (name.split('\n')[1] if '\n' in name else '', 10, C_TEXT_SUB, False)],
                           align=PP_ALIGN.CENTER)

        # 나머지 셀들
        cells = [atk, proj, tgt, special, stats]
        for j, cell_text in enumerate(cells):
            ci = j + 1
            add_rect(slide, col_x[ci], ry, col_w[ci], row_h,
                     fill_color=C_UI_BG, border_color=C_FRAME, border_width=1)
            lines = cell_text.split('\n')
            ml = [(l, 10, C_TEXT_MAIN, False) for l in lines]
            add_multiline_text(slide, col_x[ci] + 5, ry + 5, col_w[ci] - 10, row_h - 10, ml)

        # 시그니처 컬러 칩
        ci = 6
        add_rect(slide, col_x[ci], ry, col_w[ci], row_h,
                 fill_color=C_UI_BG, border_color=C_FRAME, border_width=1)
        add_rect(slide, col_x[ci] + 60, ry + 35, 80, 40,
                 fill_color=clr, border_color=C_FRAME, border_width=1)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 5: Node 특수 능력
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_nodes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DEEP)

    add_text(slide, 60, 30, 800, 40, "4. Node (마물) 특수 능력",
             font_size=22, color=C_TITLE, bold=True)
    add_rect(slide, 60, 70, 400, 2, fill_color=C_GOLD)

    # 주요 특수 Node 5종 상세
    nodes = [
        ("Shield (갑주 마물)", C_SAPPHIRE, "ShieldAbility",
         ["방어력: armor=5 (고정 차감)", "공식: Max(1, DMG - armor)", "예: Arrow Lv1(8) → 3 DMG",
          "Void Lv3+ ignoreArmor → 무시", "피드백: 'BLOCKED' + 회색 데미지"]),
        ("Regen (재생 마물)", C_EMERALD, "RegenAbility",
         ["초당 최대HP 5% 회복", "틱 간격: 0.5초 (1.25HP/틱)", "HP 50 → 순DPS 7.5 필요",
          "피드백: 초록 파티클 + HP바 글로우", "대응: 고DPS 집중"]),
        ("Phase (유령 마물)", C_PURPLE, "PhaseAbility",
         ["2초 공격가능 → 1초 무적 (반복)", "무적 중 모든 DMG 무효", "피드백: 반투명(α0.3)+보라 오라",
          "진입 0.2초 전 깜빡임 경고", "대응: Laser/Void 지속 데미지"]),
        ("Split (분열 마물)", C_FLAME, "SplitAbility",
         ["사망 시 소형 Node 2마리 분열", "분열체 HP=원본 30%, 속도 x1.2", "분열체 재분열 없음 (1회만)",
          "분열체 Bit 드롭: 2", "대응: AoE (Cannon, Lightning)"]),
        ("Boss (대마)", C_RUBY, "BossAbility + ShieldAbility",
         ["극히 높은 HP + 방어력 3", "등장: 슬로우모션 0.5초 + 경고UI", "처치: 슬로우0.5초 + 대형 폭발",
          "카메라 셰이크 + 대형 Bit 비산", "스테이지 최종 웨이브에 1체"]),
    ]

    card_w = 340
    card_h = 180
    gap = 20
    start_x = 60
    start_y = 100

    for i, (name, clr, comp, details) in enumerate(nodes):
        col = i % 3
        row = i // 3
        cx = start_x + col * (card_w + gap)
        cy = start_y + row * (card_h + gap + 30)

        add_rounded_rect(slide, cx, cy, card_w, card_h,
                         fill_color=C_UI_PANEL, border_color=clr, border_width=2)
        # 타이틀
        add_text(slide, cx + 10, cy + 8, card_w - 20, 25, name,
                 font_size=14, color=clr, bold=True)
        add_text(slide, cx + 10, cy + 30, card_w - 20, 18, comp,
                 font_size=9, color=C_TEXT_SUB)
        add_rect(slide, cx + 10, cy + 48, card_w - 20, 1, fill_color=C_FRAME)

        # 디테일
        detail_lines = [(d, 10, C_TEXT_MAIN, False) for d in details]
        add_multiline_text(slide, cx + 10, cy + 52, card_w - 20, card_h - 60, detail_lines)

    # 나머지 2개 (Boss는 이미 포함) - 빈 자리에 기본 Node들 요약
    remain_y = start_y + 2 * (card_h + gap + 30) - 30

    # 전체 요약 테이블
    add_panel(slide, 60, 520, 1440, 310, title="Node 전체 요약표")

    hdr_y2 = 570
    n_cols = [60, 200, 360, 550, 770, 1050]
    n_widths = [140, 160, 190, 220, 280, 400]
    n_headers = ["Node", "컴포넌트", "핵심 파라미터", "HP / Speed", "대응 전략", "비고"]
    for i, h in enumerate(n_headers):
        add_rect(slide, n_cols[i], hdr_y2, n_widths[i], 28,
                 fill_color=C_UI_PANEL, border_color=C_GOLD, border_width=1)
        add_text(slide, n_cols[i] + 3, hdr_y2 + 3, n_widths[i] - 6, 22, h,
                 font_size=10, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)

    n_data = [
        ("Bit (소형 마물)", "없음", "-", "20 / 1.5", "Arrow", "기본, 가장 흔함"),
        ("Quick (척후)", "없음", "speed 3.0", "12 / 3.0", "Ice 감속", "빠름"),
        ("Heavy (중갑)", "없음", "hp 80", "80 / 0.8", "화력 집중", "높은 HP"),
        ("Shield", "ShieldAbility", "armor: 5", "40 / 1.2", "고DMG or Void", "고정 차감"),
        ("Swarm (벌레)", "없음", "대량 스폰", "5 / 2.0", "AoE 필수", "극소형"),
        ("Regen", "RegenAbility", "5%HP/초", "50 / 1.0", "고DPS 집중", "회복"),
        ("Phase", "PhaseAbility", "2s/1s 사이클", "35 / 1.3", "지속 DMG", "주기적 무적"),
        ("Split", "SplitAbility", "분열2,HP30%", "30 / 1.3", "AoE", "1회 분열"),
        ("Boss", "Boss+Shield", "armor:3,고HP", "500 / 0.5", "전체 조합", "최종 웨이브"),
    ]

    rh = 22
    for i, (nm, comp, param, hp_spd, counter, note) in enumerate(n_data):
        ry = hdr_y2 + 28 + i * rh
        vals = [nm, comp, param, hp_spd, counter, note]
        for j, v in enumerate(vals):
            add_rect(slide, n_cols[j], ry, n_widths[j], rh,
                     fill_color=C_UI_BG if i % 2 == 0 else C_UI_PANEL,
                     border_color=C_FRAME, border_width=1)
            add_text(slide, n_cols[j] + 3, ry + 2, n_widths[j] - 6, rh - 4, v,
                     font_size=9, color=C_TEXT_MAIN)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 6: 보물상자 시스템
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_chest(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DEEP)

    add_text(slide, 60, 30, 800, 40, "5. 보물상자 시스템 (Hub)",
             font_size=22, color=C_TITLE, bold=True)
    add_rect(slide, 60, 70, 400, 2, fill_color=C_GOLD)

    # ── 상자 획득 조건 ──
    add_panel(slide, 60, 100, 600, 180, title="상자 획득 조건")
    acq_lines = [
        ("스테이지 첫 클리어 → 2개", 13, C_EMERALD, True),
        ("누적 웨이브 클리어 마일스톤 → 1개", 13, C_SAPPHIRE, False),
        ("매 런 종료 시 (패배 포함) → 0~1개", 13, C_TEXT_MAIN, False),
        ("  (도달 웨이브에 비례)", 11, C_TEXT_SUB, False),
        ("", 6, C_TEXT_SUB, False),
        ("상자 비용: 무료 (상자 자체가 보상)", 12, C_GOLD_BRIGHT, True),
        ("Hub에서만 열기 가능 / 인게임 중 불가", 12, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 80, 150, 560, 120, acq_lines)

    # ── 상자 내용물 ──
    add_panel(slide, 700, 100, 740, 180, title="상자 내용물 확률")
    content_lines = [
        ("75% - 타워 3후보지 (해금된 타워 풀에서 랜덤 3종, 1개 선택)", 12, C_TEXT_MAIN, False),
        ("25% - 뱃지 (타워 능력치 영구 강화 아이템)", 12, C_PURPLE, False),
        ("", 6, C_TEXT_SUB, False),
        ("후보 레벨: 모두 Lv1 (레벨업은 합성으로만)", 11, C_TEXT_SUB, False),
        ("3개 중 반드시 1개 선택 (건너뛰기 불가)", 11, C_TEXT_SUB, False),
        ("인벤토리 빈 슬롯 없으면 상자 열기 불가", 11, C_RUBY, False),
    ]
    add_multiline_text(slide, 720, 150, 700, 120, content_lines)

    # ── 오픈 연출 플로우 (도형으로) ──
    flow_y = 310
    add_text(slide, 60, flow_y, 400, 30, "상자 오픈 연출 플로우",
             font_size=16, color=C_TITLE, bold=True)
    add_rect(slide, 60, flow_y + 30, 300, 1, fill_color=C_GOLD)

    steps = [
        ("상자 등장\n0.3초 확대", C_GOLD_BRIGHT),
        ("상자 열림\n빛 방사 0.5초", C_GOLD_BRIGHT),
        ("카드 3장 등장\n(뒤집힌 상태)", C_SAPPHIRE),
        ("카드 순차 공개\n각 0.3초", C_TITLE),
        ("플레이어 선택", C_EMERALD),
        ("인벤토리로\n날아감 0.5초", C_EMERALD),
    ]
    step_w = 200
    step_h = 65
    step_gap = 30
    for i, (label, clr) in enumerate(steps):
        sx = 80 + i * (step_w + step_gap)
        sy = flow_y + 50
        # 3개씩 2줄로
        if i >= 3:
            sx = 80 + (i - 3) * (step_w + step_gap)
            sy = flow_y + 50 + step_h + 40
        add_rounded_rect(slide, sx, sy, step_w, step_h,
                         fill_color=C_UI_PANEL, border_color=clr, border_width=2)
        ls = label.split('\n')
        ml = [(ls[0], 11, clr, True)]
        if len(ls) > 1:
            ml.append((ls[1], 9, C_TEXT_SUB, False))
        add_multiline_text(slide, sx + 5, sy + 8, step_w - 10, step_h - 10, ml,
                           align=PP_ALIGN.CENTER)

    # 화살표 (1->2->3, 4->5->6)
    for i in range(2):
        ax = 80 + (i + 1) * (step_w + step_gap) - step_gap + 3
        ay = flow_y + 50 + step_h // 2
        add_rect(slide, ax, ay, step_gap - 6, 3, fill_color=C_GOLD)
    for i in range(2):
        ax = 80 + (i + 1) * (step_w + step_gap) - step_gap + 3
        ay = flow_y + 50 + step_h + 40 + step_h // 2
        add_rect(slide, ax, ay, step_gap - 6, 3, fill_color=C_GOLD)

    # ── 타워 등장 가중치 ──
    wt_y = 570
    add_panel(slide, 60, wt_y, 600, 220, title="타워 등장 가중치")

    tower_weights = [
        ("Arrow", "30", "기본, 흔함", C_TOWER_ARROW),
        ("Cannon", "25", "흔함", C_TOWER_CANNON),
        ("Ice", "20", "보통", C_TOWER_ICE),
        ("Lightning", "15", "약간 희귀", C_TOWER_LIGHTNING),
        ("Laser", "7", "희귀", C_TOWER_LASER),
        ("Void", "3", "매우 희귀", C_TOWER_VOID),
    ]
    for i, (name, weight, note, clr) in enumerate(tower_weights):
        ry = wt_y + 50 + i * 26
        add_text(slide, 80, ry, 100, 22, name, font_size=11, color=clr, bold=True)
        # 바 차트
        bar_w = int(float(weight) / 30 * 250)
        add_rect(slide, 190, ry + 5, bar_w, 14, fill_color=clr)
        add_text(slide, 190 + bar_w + 10, ry, 40, 22, weight, font_size=10, color=C_TEXT_MAIN)
        add_text(slide, 460, ry, 180, 22, note, font_size=9, color=C_TEXT_SUB)

    # ── 뱃지 시스템 ──
    add_panel(slide, 700, wt_y, 740, 220, title="뱃지 시스템")
    badge_lines = [
        ("특정 타워 타입의 능력치를 영구 강화", 12, C_PURPLE, True),
        ("인벤토리 슬롯 차지하지 않음 (별도 목록)", 11, C_TEXT_SUB, False),
        ("같은 뱃지 중첩 가능 (최대 5중첩)", 11, C_TEXT_SUB, False),
        ("", 6, C_TEXT_SUB, False),
        ("공격력 뱃지: 해당 타워 공격력 +10%/중첩", 11, C_TEXT_MAIN, False),
        ("공격속도 뱃지: 해당 타워 공속 +8%/중첩", 11, C_TEXT_MAIN, False),
        ("사거리 뱃지: 해당 타워 사거리 +0.3/중첩", 11, C_TEXT_MAIN, False),
        ("특수능력 뱃지: 해당 타워 특수 수치 강화", 11, C_TEXT_MAIN, False),
    ]
    add_multiline_text(slide, 720, wt_y + 50, 700, 160, badge_lines)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 7: 웨이브/경제
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_wave_economy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DEEP)

    add_text(slide, 60, 30, 800, 40, "6. 웨이브 규칙 / 인게임 경제",
             font_size=22, color=C_TITLE, bold=True)
    add_rect(slide, 60, 70, 400, 2, fill_color=C_GOLD)

    # ── 웨이브 규칙 ──
    add_panel(slide, 60, 100, 700, 250, title="웨이브 대기 시간")
    wave_lines = [
        ("웨이브 1 시작 전: 3초 (맵 확인 + 첫 타워 배치)", 12, C_TEXT_MAIN, False),
        ("이후 웨이브 시작 전: 3초 (기본)", 12, C_TEXT_MAIN, False),
        ("보스 웨이브 전: 5초 (총괄PD 확정)", 12, C_RUBY, True),
        ("", 6, C_TEXT_SUB, False),
        ("클리어 보너스 = 10 + (웨이브번호 × 5) Bit", 13, C_EMERALD, True),
        ("  웨이브1: 15 Bit / 웨이브2: 20 / 웨이브3: 25 / 웨이브4: 30", 11, C_TEXT_SUB, False),
        ("  × StageData.bitDropMultiplier 적용", 11, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 80, 150, 660, 190, wave_lines)

    # ── Bit 흐름 다이어그램 ──
    add_panel(slide, 800, 100, 640, 250, title="Bit 수입 / 지출 흐름")
    # 수입
    income_y = 160
    add_text(slide, 820, income_y, 100, 20, "[수입]",
             font_size=12, color=C_EMERALD, bold=True)
    incomes = ["Node 처치 드롭", "웨이브 클리어 보너스", "콤보 보너스"]
    for i, txt in enumerate(incomes):
        add_text(slide, 830, income_y + 25 + i * 20, 200, 18, "• " + txt,
                 font_size=10, color=C_EMERALD)
    # 중앙
    add_rounded_rect(slide, 1050, income_y + 10, 140, 60,
                     fill_color=C_UI_PANEL, border_color=C_GOLD_BRIGHT, border_width=2)
    add_text(slide, 1050, income_y + 25, 140, 30, "보유 Bit",
             font_size=14, color=C_GOLD_BRIGHT, bold=True, align=PP_ALIGN.CENTER)
    # 화살표
    add_rect(slide, 1030, income_y + 38, 20, 3, fill_color=C_EMERALD)
    add_rect(slide, 1190, income_y + 38, 20, 3, fill_color=C_RUBY)
    # 지출
    add_text(slide, 1220, income_y, 100, 20, "[지출]",
             font_size=12, color=C_RUBY, bold=True)
    add_text(slide, 1230, income_y + 25, 200, 18, "• 타워 구매 (placeCost)",
             font_size=10, color=C_RUBY)
    add_text(slide, 1230, income_y + 45, 200, 18, "• 판매 불가 → 환불 없음",
             font_size=10, color=C_TEXT_SUB)

    # 잔여 Bit 설명
    add_text(slide, 820, income_y + 110, 600, 40,
             "런 종료 시 잔여 Bit → 영구 저장 (Hub에서 스킬 트리 구매에 사용)",
             font_size=11, color=C_TEXT_SUB)

    # ── 타워 구매 비용표 ──
    cost_y = 370
    add_panel(slide, 60, cost_y, 500, 200, title="타워 구매 비용 (인게임 Bit)")
    tower_costs = [
        ("Arrow", "30 Bit", C_TOWER_ARROW),
        ("Cannon", "60 Bit", C_TOWER_CANNON),
        ("Ice", "45 Bit", C_TOWER_ICE),
        ("Lightning", "80 Bit", C_TOWER_LIGHTNING),
        ("Laser", "120 Bit", C_TOWER_LASER),
        ("Void", "150 Bit", C_TOWER_VOID),
    ]
    for i, (name, cost, clr) in enumerate(tower_costs):
        ry = cost_y + 50 + i * 24
        add_text(slide, 80, ry, 120, 20, name, font_size=11, color=clr, bold=True)
        add_text(slide, 230, ry, 100, 20, cost, font_size=11, color=C_TEXT_MAIN)

    # ── 콤보 시스템 ──
    add_panel(slide, 600, cost_y, 840, 200, title="콤보 보너스 시스템")
    combo_lines = [
        ("3초 이내 연속 처치 → 콤보 카운터 증가", 12, C_TEXT_MAIN, False),
        ("", 4, C_TEXT_SUB, False),
        ("1~4 콤보: 보너스 없음, 카운터만 표시", 11, C_TEXT_SUB, False),
        ("5~9 콤보: +50% Bit, 'COMBO x5!' + 파티클 확대", 11, C_FLAME, False),
        ("10~19 콤보: +100% Bit, 'COMBO x10!' + 화면 흔들림", 11, C_GOLD_BRIGHT, False),
        ("20+ 콤보: +200% Bit, 'MEGA COMBO!' + 화면 플래시", 11, C_RUBY, True),
    ]
    add_multiline_text(slide, 620, cost_y + 50, 800, 140, combo_lines)

    # ── 경제 시뮬레이션 ──
    sim_y = 590
    add_panel(slide, 60, sim_y, 1380, 290, title="경제 밸런스 시뮬레이션 (스테이지 1, 영구 업그레이드 0)")
    sim_lines = [
        ("런 시작: 0 Bit (시작 Bit 스킬 없음)", 11, C_TEXT_SUB, False),
        ("", 4, C_TEXT_SUB, False),
        ("웨이브 1 (Bit x5): 처치 15 + 클리어 15 = 30 Bit → Arrow 1기 구매 가능", 11, C_TEXT_MAIN, False),
        ("웨이브 2 (Bit x8): 처치 24 + 클리어 20 = 44 Bit → 누적 ~74 Bit", 11, C_TEXT_MAIN, False),
        ("웨이브 3 (Bit x6 + Quick x3): 처치 30 + 클리어 25 = 55 Bit → 대부분 여기서 사망", 11, C_RUBY, False),
        ("", 4, C_TEXT_SUB, False),
        ("매 런 영구 업그레이드: 공격력 Lv1(50), 공속 Lv1(80) 등 1~2개 가능", 11, C_EMERALD, False),
        ("", 4, C_TEXT_SUB, False),
        ("총 가능 Bit (전부 처치 + 클리어): 179 + 90 = 269 Bit", 12, C_GOLD_BRIGHT, True),
        ("스테이지 클리어 시: +Core 2", 12, C_PURPLE, True),
        ("", 4, C_TEXT_SUB, False),
        ("첫 클리어 예상: 런 11~20회, 총 플레이 시간 ~20~30분", 11, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 80, sim_y + 48, 1340, 230, sim_lines)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 8: 전투 피드백
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_feedback(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DEEP)

    add_text(slide, 60, 30, 800, 40, "7. 전투 피드백 / 이펙트 트리거",
             font_size=22, color=C_TITLE, bold=True)
    add_rect(slide, 60, 70, 400, 2, fill_color=C_GOLD)

    # 피드백 테이블
    hdr_y = 95
    f_cols = [60, 290, 660, 960, 1140]
    f_widths = [230, 370, 300, 180, 320]
    f_headers = ["이벤트", "시각 피드백", "사운드", "카메라", "조건"]
    for i, h in enumerate(f_headers):
        add_rect(slide, f_cols[i], hdr_y, f_widths[i], 28,
                 fill_color=C_UI_PANEL, border_color=C_GOLD, border_width=1)
        add_text(slide, f_cols[i] + 3, hdr_y + 3, f_widths[i] - 6, 22, h,
                 font_size=10, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)

    fb_data = [
        ("타워 공격 발사", "발사 애니메이션 (2f)", "타워별 공격음", "-", "매 공격"),
        ("투사체 적중", "타격 스파크 (타워 컬러)", "히트음 (짧고 경쾌)", "-", "매 적중"),
        ("Node 처치", "파괴 파티클 + '+N Bit' 팝업", "처치음 (타입별)", "-", "매 처치"),
        ("동시 다수 처치 (AoE)", "큰 폭발 + 합산 Bit 팝업", "폭발음 (강조)", "미세 흔들림", "3마리+ 동시"),
        ("콤보 5+", "'COMBO xN!' + 파티클 확대", "콤보 효과음 (상승)", "-", "콤보 5+"),
        ("Node 기지 도착", "빨간 플래시 + HP바 감소", "피격음", "테두리 붉게", "HP 감소 시"),
        ("웨이브 시작", "'Wave N' 배너 슬라이드인", "경고음 (짧게)", "-", "매 웨이브"),
        ("웨이브 클리어", "'WAVE CLEAR' + 가장자리 플래시", "클리어 팡파르", "-", "매 클리어"),
        ("보스 등장", "'WARNING!' + 화면 어두워짐", "보스 등장 사운드", "슬로우 0.5초", "보스 웨이브"),
        ("보스 처치", "대형 폭발 + Bit 비산", "보스 처치 팡파르", "슬로우 + 흔들림", "보스 사망"),
        ("타워 배치", "마법진 출현 이펙트", "배치음 (돌)", "-", "배치 시"),
        ("타워 합성", "발광 + 'Lv N!' 팝업", "강화음 (글리산도)", "Lv4시 흔들림", "합성 시"),
        ("스테이지 클리어", "밝기↑ + Core등장 + 금빛파티클", "대형 팡파르", "-", "스테이지 클리어"),
        ("HP < 30%", "붉은 바이넷 (지속)", "심장박동 (반복)", "-", "지속 경고"),
        ("HP < 10%", "바이넷 강도↑ + 펄스", "심장박동 빨라짐", "-", "위급 경고"),
    ]

    rh = 22
    for i, (evt, visual, sound, cam, cond) in enumerate(fb_data):
        ry = hdr_y + 28 + i * rh
        vals = [evt, visual, sound, cam, cond]
        for j, v in enumerate(vals):
            bg = C_UI_BG if i % 2 == 0 else C_UI_PANEL
            add_rect(slide, f_cols[j], ry, f_widths[j], rh,
                     fill_color=bg, border_color=C_FRAME, border_width=1)
            add_text(slide, f_cols[j] + 3, ry + 2, f_widths[j] - 6, rh - 4, v,
                     font_size=8, color=C_TEXT_MAIN)

    # 카메라 셰이크 명세
    shake_y = hdr_y + 28 + len(fb_data) * rh + 20
    add_panel(slide, 60, shake_y, 700, 160, title="카메라 셰이크 명세")
    shake_lines = [
        ("Cannon AoE (3+타격): intensity=0.05, duration=0.1s", 11, C_FLAME, False),
        ("보스 등장: intensity=0.08, duration=0.3s", 11, C_RUBY, False),
        ("보스 처치: intensity=0.15, duration=0.5s", 11, C_RUBY, True),
        ("타워 Lv4(MAX) 합성: intensity=0.03, duration=0.1s", 11, C_GOLD_BRIGHT, False),
        ("기지 큰 피해 (3+): intensity=0.04, duration=0.15s", 11, C_TEXT_MAIN, False),
    ]
    add_multiline_text(slide, 80, shake_y + 48, 660, 100, shake_lines)

    # 크리티컬 시스템
    add_panel(slide, 800, shake_y, 640, 160, title="크리티컬 시스템")
    crit_lines = [
        ("스킬 트리 '크리티컬 해금' 노드 필요", 11, C_TEXT_SUB, False),
        ("기본 확률: 15% / 기본 배율: 2.0x", 12, C_GOLD_BRIGHT, True),
        ("시각: 금색 데미지 숫자 + 별 파티클 + 잠깐 밝아짐", 11, C_TEXT_MAIN, False),
        ("사운드: 강조된 '크랭' 사운드", 11, C_TEXT_MAIN, False),
    ]
    add_multiline_text(slide, 820, shake_y + 48, 600, 100, crit_lines)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 9: 밸런싱 - 스테이지 1 웨이브 구성
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_balancing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DEEP)

    add_text(slide, 60, 30, 800, 40, "8. 밸런싱 - 스테이지 1 웨이브 구성",
             font_size=22, color=C_TITLE, bold=True)
    add_rect(slide, 60, 70, 400, 2, fill_color=C_GOLD)

    # 웨이브 테이블
    hdr_y = 100
    w_cols = [60, 160, 440, 570, 660, 780, 1020]
    w_widths = [100, 280, 130, 90, 120, 240, 440]
    w_headers = ["웨이브", "Node 구성", "스폰 간격", "지연", "총 Node", "총 Bit (처치)", "설계 의도"]
    for i, h in enumerate(w_headers):
        add_rect(slide, w_cols[i], hdr_y, w_widths[i], 30,
                 fill_color=C_UI_PANEL, border_color=C_GOLD, border_width=1)
        add_text(slide, w_cols[i] + 3, hdr_y + 4, w_widths[i] - 6, 22, h,
                 font_size=10, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)

    waves = [
        ("1", "Bit x5", "1.2초", "3.0초", "5", "15", "학습. 타워가 잡는다."),
        ("2", "Bit x8", "1.0초", "3.0초", "8", "24", "물량 증가. 타워 추가 유도."),
        ("3", "Bit x6 + Quick x3", "0.8초", "3.0초", "9", "30", "빠른 적 등장. 긴장감."),
        ("4", "Bit x10 + Quick x5", "0.7초", "3.0초", "15", "50", "물량 압박. 초기 사망 구간."),
        ("5", "Bit x8 + Quick x4 + Heavy x2", "0.6초", "5.0초", "14", "60", "최종. 첫 클리어 시 성취감."),
    ]

    rh = 50
    for i, (wave, comp, interval, delay, total, bits, intent) in enumerate(waves):
        ry = hdr_y + 30 + i * rh
        vals = [wave, comp, interval, delay, total, bits, intent]
        row_color = C_RUBY if i == 4 else C_FLAME if i == 3 else C_TEXT_MAIN
        for j, v in enumerate(vals):
            bg = C_UI_BG if i % 2 == 0 else C_UI_PANEL
            add_rect(slide, w_cols[j], ry, w_widths[j], rh,
                     fill_color=bg, border_color=C_FRAME, border_width=1)
            clr = row_color if j == 6 else C_TEXT_MAIN
            add_text(slide, w_cols[j] + 5, ry + 12, w_widths[j] - 10, rh - 10, v,
                     font_size=10, color=clr)

    # 합계
    sum_y = hdr_y + 30 + 5 * rh + 10
    add_panel(slide, 60, sum_y, 700, 90, title="합계")
    sum_lines = [
        ("처치 Bit: 15+24+30+50+60 = 179 Bit", 12, C_EMERALD, False),
        ("클리어 보너스: 15+20+25+30 = 90 Bit | 합계: 269 Bit + Core 2", 12, C_GOLD_BRIGHT, True),
    ]
    add_multiline_text(slide, 80, sum_y + 48, 660, 40, sum_lines)

    # DPS 검증
    dps_y = sum_y + 110
    add_panel(slide, 60, dps_y, 700, 180, title="DPS 비율 검증 (영구 업그레이드 0)")
    dps_lines = [
        ("Arrow Lv1 DPS: 10.0", 12, C_TOWER_ARROW, True),
        ("Bit Node (HP 20): 처치 2.0초 (목표 1.5~2.5초 ✓)", 11, C_TEXT_MAIN, False),
        ("Quick Node (HP 12, 속도 3.0): 처치 1.2초, 통과 가능성 있음 (의도)", 11, C_TEXT_MAIN, False),
        ("Heavy Node (HP 80): Arrow 1기로 8.0초 (목표 6~10초 ✓)", 11, C_TEXT_MAIN, False),
        ("", 4, C_TEXT_SUB, False),
        ("런당 생존 시간 (초기): 30초~1분 → 빠른 죽음 → 빠른 재시도", 11, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 80, dps_y + 48, 660, 120, dps_lines)

    # 성장 곡선
    grow_y = sum_y
    add_panel(slide, 800, grow_y, 640, 290, title="첫 런 ~ 첫 클리어 경험 곡선")
    grow_lines = [
        ("런 1 (업그레이드 0):", 12, C_TEXT_MAIN, True),
        ("  Arrow 1기 자동배치 → W1 여유 → W2 빠듯 → W3~4 사망", 10, C_TEXT_SUB, False),
        ("  획득: ~80~100 Bit", 10, C_EMERALD, False),
        ("", 4, C_TEXT_SUB, False),
        ("런 2~5 (공격력+1, 공속+1):", 12, C_TEXT_MAIN, True),
        ("  영구 업그레이드 체감. W3~4까지 안정 도달", 10, C_TEXT_SUB, False),
        ("  획득: ~100~150 Bit / 런", 10, C_EMERALD, False),
        ("", 4, C_TEXT_SUB, False),
        ("런 6~10 (공격+3, 공속+2, HP+2):", 12, C_TEXT_MAIN, True),
        ("  DPS 체감 상승. 타워 2~3기 운용. W4~5 진입", 10, C_TEXT_SUB, False),
        ("", 4, C_TEXT_SUB, False),
        ("런 11~20: 스테이지 1 클리어 가능", 12, C_GOLD_BRIGHT, True),
        ("총 플레이 시간: ~20~30분", 12, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 820, grow_y + 48, 600, 230, grow_lines)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 10: 구현 우선순위
# ═══════════════════════════════════════════════════════════════════════════
def build_slide_priority(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DEEP)

    add_text(slide, 60, 30, 800, 40, "9. 구현 우선순위",
             font_size=22, color=C_TITLE, bold=True)
    add_rect(slide, 60, 70, 400, 2, fill_color=C_GOLD)

    phases = [
        ("Phase 1: 핵심 전투 루프", C_EMERALD, [
            ("PlacementGrid (격자 배치 판정)", "중"),
            ("TowerDragController (드래그 배치)", "중"),
            ("타워 인벤토리 시스템", "중"),
            ("웨이브 클리어 이벤트 + 보너스", "하"),
            ("콤보 시스템", "하"),
            ("데미지 숫자 팝업", "하"),
        ]),
        ("Phase 2: 타워 특수 능력", C_FLAME, [
            ("CannonProjectile (AoE)", "중"),
            ("Node 감속 + IceProjectile", "중"),
            ("LightningTower (체인 공격)", "상"),
            ("LaserTower (빔 + 쿨다운)", "상"),
            ("VoidTower (범위 지속)", "중"),
        ]),
        ("Phase 3: Node 특수 능력", C_PURPLE, [
            ("NodeAbility 기반 구조 변경", "중"),
            ("ShieldAbility (방어력)", "하"),
            ("RegenAbility (회복)", "하"),
            ("PhaseAbility (무적)", "중"),
            ("SplitAbility (분열)", "상"),
            ("BossAbility (연출)", "중"),
        ]),
        ("Phase 4: 시스템 보완", C_SAPPHIRE, [
            ("타워 배치 SaveData 저장/복원", "중"),
            ("RunModifiers 확장", "하"),
            ("크리티컬 시스템", "하"),
            ("다중 경로 지원", "중"),
            ("전투 피드백/이펙트", "중~상"),
            ("보물상자 시스템 (Hub)", "중"),
        ]),
    ]

    card_w = 340
    card_gap = 15
    start_x = 60
    start_y = 100

    for pi, (phase_title, phase_color, tasks) in enumerate(phases):
        cx = start_x + pi * (card_w + card_gap)
        cy = start_y

        # Phase 카드
        card_h = 60 + len(tasks) * 28 + 20
        add_rounded_rect(slide, cx, cy, card_w, card_h,
                         fill_color=C_UI_PANEL, border_color=phase_color, border_width=2)

        # Phase 제목
        add_text(slide, cx + 10, cy + 10, card_w - 20, 30, phase_title,
                 font_size=15, color=phase_color, bold=True)
        add_rect(slide, cx + 10, cy + 42, card_w - 20, 1, fill_color=C_FRAME)

        # 태스크 목록
        for ti, (task, diff) in enumerate(tasks):
            ty = cy + 52 + ti * 28
            # 난이도 컬러
            diff_color = C_RUBY if '상' in diff else C_FLAME if '중' in diff else C_EMERALD
            # 순번
            add_text(slide, cx + 10, ty, 25, 22,
                     f"{pi+1}-{ti+1}", font_size=9, color=C_TEXT_SUB)
            # 태스크명
            add_text(slide, cx + 38, ty, card_w - 100, 22,
                     task, font_size=10, color=C_TEXT_MAIN)
            # 난이도 태그
            add_rounded_rect(slide, cx + card_w - 55, ty + 2, 40, 18,
                             fill_color=C_UI_BG, border_color=diff_color, border_width=1)
            add_text(slide, cx + card_w - 55, ty + 3, 40, 16,
                     diff, font_size=8, color=diff_color, align=PP_ALIGN.CENTER)

    # 하단 배속/일시정지 참고
    note_y = 730
    add_panel(slide, 60, note_y, 1380, 150, title="배속 / 일시정지 시스템 (보완)")
    note_lines = [
        ("배속 단계: x1, x2, x3 → x1만 기본, x2/x3은 '배속 해금' Core 노드 필요", 12, C_TEXT_MAIN, False),
        ("일시정지 중 타워 배치/합성: 불가 (총괄PD 확정)", 12, C_RUBY, True),
        ("배속 중 이펙트: Time.deltaTime 기반이므로 자동 단축", 12, C_TEXT_SUB, False),
        ("", 4, C_TEXT_SUB, False),
        ("스테이지별 권장: S1 Arrow x3~4 / S2 +Cannon / S3 +Ice / S4~5 +Lightning", 11, C_TEXT_SUB, False),
        ("S6~7 +Laser / S8~9 +Void / S10 전 타워 조합", 11, C_TEXT_SUB, False),
    ]
    add_multiline_text(slide, 80, note_y + 48, 1340, 90, note_lines)


# ═══════════════════════════════════════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_cover(prs)           # 1. 표지
    build_slide_run_flow(prs)        # 2. 런 플로우
    build_slide_tower_placement(prs) # 3. 타워 배치/합성
    build_slide_towers(prs)          # 4. 타워 6종
    build_slide_nodes(prs)           # 5. Node 특수능력
    build_slide_chest(prs)           # 6. 보물상자
    build_slide_wave_economy(prs)    # 7. 웨이브/경제
    build_slide_feedback(prs)        # 8. 전투 피드백
    build_slide_balancing(prs)       # 9. 밸런싱
    build_slide_priority(prs)        # 10. 구현 우선순위

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"생성 완료: {OUT_PATH}")
    print(f"슬라이드 수: {len(prs.slides)}")
    print(f"파일 크기: {os.path.getsize(OUT_PATH) / 1024:.1f} KB")


if __name__ == "__main__":
    main()
