#!/usr/bin/env python3
"""
Soulspire TD - UI 레이아웃 명세서 PPT 생성기
해상도: 1920x1080 (16:9)
모든 좌표는 픽셀 단위이며, PPT용 EMU로 변환됩니다.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── 상수 ─────────────────────────────────────────────────────────────────
OUT_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "UI_Spec_v0.1.pptx",
)

# 1 px = 9525 EMU (96 DPI 기준, Windows 표준 DPI)
PX = 9525

# 슬라이드 크기 (1920x1080 px)
SLIDE_W = 1920 * PX  # 18288000 EMU
SLIDE_H = 1080 * PX  # 10287000 EMU

# ─── 색상 팔레트 (네온 회로판 테마) ───────────────────────────────────────
C_BG_DARK = RGBColor(0x0B, 0x0F, 0x1A)       # #0B0F1A - 메인 배경
C_PANEL = RGBColor(0x12, 0x1A, 0x2A)          # #121A2A - 패널 배경
C_PANEL_LIGHT = RGBColor(0x1A, 0x24, 0x3A)    # #1A243A - 밝은 패널
C_NEON_GREEN = RGBColor(0x2B, 0xFF, 0x88)     # #2BFF88 - 기본 강조
C_NEON_BLUE = RGBColor(0x37, 0xB6, 0xFF)      # #37B6FF - 보조 강조
C_NEON_PURPLE = RGBColor(0xB8, 0x6C, 0xFF)    # #B86CFF - 3차 강조
C_ORANGE = RGBColor(0xFF, 0x9A, 0x3D)         # #FF9A3D - 경고/하이라이트
C_RED = RGBColor(0xFF, 0x4D, 0x5A)            # #FF4D5A - 경보/위험
C_YELLOW = RGBColor(0xFF, 0xD8, 0x4D)         # #FFD84D - 희귀/골드
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT_MAIN = RGBColor(0xD8, 0xE4, 0xFF)      # #D8E4FF - 기본 텍스트
C_TEXT_DIM = RGBColor(0xAF, 0xC3, 0xE8)       # #AFC3E8 - 보조 텍스트
C_BORDER = RGBColor(0x5B, 0x6B, 0x8A)         # #5B6B8A - 기본 테두리
C_HP_GREEN = RGBColor(0x44, 0xCC, 0x44)       # #44CC44 - HP 바 초록
C_HP_RED = RGBColor(0xCC, 0x22, 0x22)         # #CC2222 - HP 바 빈칸


def px(val):
    """픽셀 값을 EMU로 변환."""
    return int(val * PX)


def add_bg(slide, color=C_BG_DARK):
    """슬라이드 배경을 단색으로 채움."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=None, border_color=None, border_width=1):
    """사각형 도형 추가."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, x, y, w, h, fill_color=None, border_color=None, border_width=1):
    """둥근 사각형 도형 추가."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, font_size=14, color=C_TEXT_MAIN,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """텍스트 상자 추가."""
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    tf.auto_size = None
    # 세로 정렬 설정
    txBox.text_frame.paragraphs[0].alignment = align
    return txBox


def add_multiline_text(slide, x, y, w, h, lines, font_size=12, color=C_TEXT_MAIN,
                       bold=False, align=PP_ALIGN.LEFT):
    """여러 줄 텍스트 상자 추가 (텍스트, 크기, 색상, 굵기) 튜플 리스트."""
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            txt, sz, clr, b = line_data, font_size, color, bold
        else:
            txt = line_data[0]
            sz = line_data[1] if len(line_data) > 1 else font_size
            clr = line_data[2] if len(line_data) > 2 else color
            b = line_data[3] if len(line_data) > 3 else bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = b
        p.alignment = align
    return txBox


def add_spec_label(slide, x, y, text, font_size=8, color=C_TEXT_DIM):
    """작은 명세 주석 라벨 추가."""
    return add_text(slide, x, y, 300, 20, text, font_size=font_size, color=color)


def add_dimension_label(slide, x, y, w, h, text):
    """치수/위치 주석 추가."""
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7)
    p.font.color.rgb = C_ORANGE
    p.font.bold = False
    p.alignment = PP_ALIGN.LEFT
    return txBox


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 1: 타이틀 화면
# ═══════════════════════════════════════════════════════════════════════════════
def build_slide_title_screen(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드
    add_bg(slide, C_BG_DARK)

    # 슬라이드 제목 (명세서 제목, 게임 레이아웃 바깥)
    add_text(slide, 20, 10, 500, 30, "슬라이드 1: 타이틀 화면 (게임 시작)",
             font_size=10, color=C_ORANGE, bold=True)

    # ── 배경 (전체 화면) ──
    add_rect(slide, 0, 40, 1920, 1040, fill_color=C_BG_DARK, border_color=C_BORDER)
    add_dimension_label(slide, 5, 42, 300, 20,
                        "배경: (0,0) 1920x1080 #0B0F1A")

    # ── 게임 제목 ──
    title_x, title_y, title_w, title_h = 460, 200, 1000, 100
    add_rect(slide, title_x, title_y, title_w, title_h,
             fill_color=None, border_color=C_NEON_GREEN, border_width=1)
    add_text(slide, title_x, title_y + 10, title_w, title_h,
             "SOULSPIRE TD", font_size=48, color=C_NEON_GREEN,
             bold=True, align=PP_ALIGN.CENTER)
    add_dimension_label(slide, title_x, title_y - 18, 500, 20,
                        f"제목: ({title_x},{title_y-40}) {title_w}x{title_h}px, 48pt 굵게, #2BFF88, Anchor=Center-Center")

    # ── 부제목 ──
    sub_y = 320
    add_text(slide, 460, sub_y, 1000, 40,
             "네트워크를 지켜라. Node를 격파하라.",
             font_size=16, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
    add_dimension_label(slide, 460, sub_y - 16, 500, 20,
                        f"부제목: (460,{sub_y}) 1000x40px, 16pt, #AFC3E8")

    # ── 버튼 (중앙 정렬, 세로 배치) ──
    btn_x, btn_w, btn_h = 760, 400, 60
    btn_gap = 20
    buttons = [
        ("시작", 480, C_NEON_GREEN, C_BG_DARK),
        ("설정", 560, C_NEON_BLUE, C_BG_DARK),
        ("종료", 640, C_RED, C_BG_DARK),
    ]
    for label, btn_y, border_c, fill_c in buttons:
        add_rounded_rect(slide, btn_x, btn_y, btn_w, btn_h,
                         fill_color=C_PANEL, border_color=border_c, border_width=2)
        add_text(slide, btn_x, btn_y + 12, btn_w, btn_h - 12,
                 f"[ {label} ]", font_size=20, color=border_c,
                 bold=True, align=PP_ALIGN.CENTER)
        add_dimension_label(slide, btn_x + btn_w + 10, btn_y + 10, 500, 20,
                            f"버튼 '{label}': ({btn_x},{btn_y}) {btn_w}x{btn_h}px, 테두리 #{border_c}, 20pt")

    # ── 명세 노트 ──
    notes_data = [
        ("명세 노트 - 타이틀 화면", 11, C_ORANGE, True),
        ("", 8, C_TEXT_DIM, False),
        ("레이아웃:", 10, C_TEXT_MAIN, True),
        ("- 해상도: 1920x1080", 9, C_TEXT_DIM, False),
        ("- 모든 요소 가로 중앙 정렬 (Anchor: Center)", 9, C_TEXT_DIM, False),
        ("- 제목 Y=200, 버튼 시작 Y=480, 간격=20px", 9, C_TEXT_DIM, False),
        ("", 8, C_TEXT_DIM, False),
        ("Unity RectTransform:", 10, C_TEXT_MAIN, True),
        ("- 배경: Stretch-Stretch, 모든 오프셋=0", 9, C_TEXT_DIM, False),
        ("- 제목: Anchor=Top-Center, Pivot=(0.5, 0.5)", 9, C_TEXT_DIM, False),
        ("- 버튼: Anchor=Top-Center, Pivot=(0.5, 0.5)", 9, C_TEXT_DIM, False),
        ("- 버튼 간격: 중심 간 80px", 9, C_TEXT_DIM, False),
        ("", 8, C_TEXT_DIM, False),
        ("색상:", 10, C_TEXT_MAIN, True),
        ("- 배경: #0B0F1A | 패널: #121A2A", 9, C_TEXT_DIM, False),
        ("- 시작 버튼 테두리: #2BFF88 | 설정: #37B6FF | 종료: #FF4D5A", 9, C_TEXT_DIM, False),
        ("- 제목 텍스트: #2BFF88 48pt 굵게", 9, C_TEXT_DIM, False),
        ("", 8, C_TEXT_DIM, False),
        ("인터랙션:", 10, C_TEXT_MAIN, True),
        ("- 호버: 테두리 발광 강도 +50%, 채우기 #1A243A로 밝게", 9, C_TEXT_DIM, False),
        ("- 클릭: 0.1초간 0.95배 축소", 9, C_TEXT_DIM, False),
    ]
    add_multiline_text(slide, 20, 780, 600, 280, notes_data)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 2: Hub 화면 (스킬 트리)
# ═══════════════════════════════════════════════════════════════════════════════
def build_slide_hub(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DARK)

    add_text(slide, 20, 10, 500, 30, "슬라이드 2: Hub 화면 (스킬 트리)",
             font_size=10, color=C_ORANGE, bold=True)

    # ── 상단 바 ──
    top_bar_h = 50
    add_rect(slide, 0, 40, 1920, top_bar_h, fill_color=C_PANEL, border_color=C_BORDER)
    add_dimension_label(slide, 5, 25, 500, 20,
                        f"상단 바: (0,0) 1920x{top_bar_h}px, 채우기=#121A2A, 하단 테두리=#5B6B8A")

    # Bit 카운터
    add_text(slide, 20, 48, 200, 35, "Bit: 1,250", font_size=16, color=C_NEON_GREEN, bold=True)
    add_dimension_label(slide, 20, 85, 400, 20,
                        "Bit: (20,8) 200x35px, 16pt 굵게, #2BFF88, Anchor=Left")

    # Core 카운터
    add_text(slide, 240, 48, 200, 35, "Core: 5", font_size=16, color=C_NEON_PURPLE, bold=True)
    add_dimension_label(slide, 240, 85, 400, 20,
                        "Core: (240,8) 200x35px, 16pt 굵게, #B86CFF")

    # 방치 Bit 알림 (우측 정렬)
    idle_x, idle_w = 1500, 400
    add_rounded_rect(slide, idle_x, 46, idle_w, 38,
                     fill_color=C_PANEL_LIGHT, border_color=C_YELLOW, border_width=1)
    add_text(slide, idle_x + 10, 50, idle_w - 20, 30,
             "방치 Bit: 150  [수령!]", font_size=13, color=C_YELLOW, bold=True)
    add_dimension_label(slide, idle_x, 85, 500, 20,
                        f"방치 버튼: ({idle_x},6) {idle_w}x38px, #FFD84D 테두리, 방치 Bit>0일 때 표시")

    # ── 스킬 트리 영역 (메인 콘텐츠) ──
    tree_y = 95
    tree_h = 890
    add_rect(slide, 0, tree_y, 1920, tree_h,
             fill_color=C_BG_DARK, border_color=C_BORDER)
    add_dimension_label(slide, 5, tree_y + 2, 600, 20,
                        f"스킬 트리 영역: (0,{tree_y-40}) 1920x{tree_h}px, 스크롤+줌 가능, Anchor=Stretch")

    # 중앙 노드
    cx, cy = 960, 540
    node_size = 70
    add_rounded_rect(slide, cx - node_size // 2, cy - node_size // 2,
                     node_size, node_size,
                     fill_color=C_NEON_GREEN, border_color=C_WHITE, border_width=2)
    add_text(slide, cx - node_size // 2, cy - 10, node_size, 30,
             "CORE", font_size=11, color=C_BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_dimension_label(slide, cx + 40, cy - 40, 400, 20,
                        f"Core 노드: 화면 중앙, {node_size}x{node_size}px, 채우기=#2BFF88")

    # 샘플 트리 노드 (중앙에서 4방향)
    node_w, node_h = 60, 60
    tree_nodes = [
        # (오프셋_x, 오프셋_y, 라벨, 상태, 색상)
        (0, -120, "ATK+", "active", C_NEON_GREEN),
        (0, 120, "HP+", "active", C_NEON_GREEN),
        (-140, 0, "RNG+", "purchasable", C_NEON_BLUE),
        (140, 0, "SPD+", "purchasable", C_NEON_BLUE),
        (0, -240, "Bit+", "locked", C_BORDER),
        (140, -120, "CRIT", "locked", C_BORDER),
        (-140, -120, "HEAL", "hidden", C_PANEL),
        (280, 0, "Cannon", "locked", C_BORDER),
        (-280, 0, "Idle", "hidden", C_PANEL),
        (0, 240, "Start$", "locked", C_BORDER),
        (140, 120, "Ice", "hidden", C_PANEL),
    ]
    for ox, oy, label, state, color in tree_nodes:
        nx = cx + ox - node_w // 2
        ny = cy + oy - node_h // 2
        bw = 2 if state in ("active", "purchasable") else 1
        fill = C_PANEL if state != "active" else RGBColor(0x15, 0x30, 0x20)
        add_rounded_rect(slide, nx, ny, node_w, node_h,
                         fill_color=fill, border_color=color, border_width=bw)
        add_text(slide, nx, ny + 15, node_w, 30,
                 label, font_size=9, color=color, bold=(state == "active"),
                 align=PP_ALIGN.CENTER)

    # 연결선 (얇은 사각형으로 단순화)
    line_color = C_NEON_GREEN
    # 중앙에서 위로 세로선
    add_rect(slide, cx - 1, cy - node_size // 2 - 55, 3, 55,
             fill_color=line_color)
    # 중앙에서 아래로 세로선
    add_rect(slide, cx - 1, cy + node_size // 2, 3, 55,
             fill_color=line_color)

    # 노드 상태 범례
    legend_x, legend_y = 30, tree_y + 10
    add_rect(slide, legend_x, legend_y, 280, 160,
             fill_color=C_PANEL, border_color=C_BORDER)
    legend_lines = [
        ("노드 상태:", 10, C_TEXT_MAIN, True),
        ("활성: 밝은 채우기 + 발광 테두리", 9, C_NEON_GREEN, False),
        ("구매 가능: 깜빡이는 테두리 + 비용 표시", 9, C_NEON_BLUE, False),
        ("잠김: 흐릿함, 선행 조건 필요", 9, C_BORDER, False),
        ("숨김: 인접 노드 해금 전까지 보이지 않음", 9, C_PANEL_LIGHT, False),
        ("", 8, C_TEXT_DIM, False),
        ("노드 크기: 60x60px, 간격: 80px", 8, C_TEXT_DIM, False),
        ("연결선: 2px, 노드와 동일 색상", 8, C_TEXT_DIM, False),
    ]
    add_multiline_text(slide, legend_x + 10, legend_y + 8, 260, 150, legend_lines)

    # ── 하단 바 ──
    bot_bar_y = 990
    bot_bar_h = 55
    add_rect(slide, 0, bot_bar_y, 1920, bot_bar_h,
             fill_color=C_PANEL, border_color=C_BORDER)
    add_dimension_label(slide, 5, bot_bar_y - 16, 600, 20,
                        f"하단 바: (0,{bot_bar_y - 40}) 1920x{bot_bar_h}px, 채우기=#121A2A")

    # 스테이지 선택 드롭다운
    dd_x, dd_w = 20, 300
    add_rounded_rect(slide, dd_x, bot_bar_y + 8, dd_w, 38,
                     fill_color=C_PANEL_LIGHT, border_color=C_NEON_BLUE)
    add_text(slide, dd_x + 10, bot_bar_y + 14, dd_w - 40, 30,
             "스테이지 1: 데이터 스트림  v", font_size=12, color=C_NEON_BLUE)
    add_dimension_label(slide, dd_x, bot_bar_y + 48, 500, 20,
                        f"스테이지 드롭다운: ({dd_x},{bot_bar_y - 40 + 8}) {dd_w}x38px")

    # 출격 버튼
    deploy_x, deploy_w = 810, 300
    add_rounded_rect(slide, deploy_x, bot_bar_y + 5, deploy_w, 44,
                     fill_color=RGBColor(0x15, 0x40, 0x25), border_color=C_NEON_GREEN, border_width=2)
    add_text(slide, deploy_x, bot_bar_y + 12, deploy_w, 35,
             "[ 출격 ]", font_size=18, color=C_NEON_GREEN,
             bold=True, align=PP_ALIGN.CENTER)
    add_dimension_label(slide, deploy_x, bot_bar_y + 52, 400, 20,
                        f"출격 버튼: ({deploy_x},{bot_bar_y - 40 + 5}) {deploy_w}x44px, #2BFF88")

    # 설정 + 종료 버튼 (우측)
    for i, (label, clr) in enumerate([("설정", C_NEON_BLUE), ("종료", C_RED)]):
        bx = 1700 + i * 110
        add_rounded_rect(slide, bx, bot_bar_y + 8, 100, 38,
                         fill_color=C_PANEL_LIGHT, border_color=clr)
        add_text(slide, bx, bot_bar_y + 14, 100, 30,
                 label, font_size=11, color=clr, align=PP_ALIGN.CENTER)
    add_dimension_label(slide, 1700, bot_bar_y + 48, 400, 20,
                        "설정/종료: 각 100x38px, 우측 정렬, 간격=10px")

    # 명세 노트
    notes_data = [
        ("명세 노트 - Hub 화면", 11, C_ORANGE, True),
        ("", 6, C_TEXT_DIM, False),
        ("구조: 상단 바 (50px) + 스킬 트리 (나머지) + 하단 바 (55px)", 9, C_TEXT_DIM, False),
        ("스킬 트리: 이동=우클릭 드래그, 줌=스크롤 휠", 9, C_TEXT_DIM, False),
        ("노드 클릭 -> 구매 팝업 (슬라이드 5 참조)", 9, C_TEXT_DIM, False),
        ("진행도 표시: 트리 영역 우측 상단에 '해금: 12/45'", 9, C_TEXT_DIM, False),
        ("Anchor: 상단 바=Top-Stretch, 하단 바=Bottom-Stretch, 트리=Stretch-All", 9, C_TEXT_DIM, False),
    ]
    add_multiline_text(slide, 1350, tree_y + 10, 550, 140, notes_data)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 3: 인게임 화면
# ═══════════════════════════════════════════════════════════════════════════════
def build_slide_ingame(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DARK)

    add_text(slide, 20, 10, 500, 30, "슬라이드 3: 인게임 화면 (런 진행 중)",
             font_size=10, color=C_ORANGE, bold=True)

    # ── 상단 HUD 바 ──
    hud_y = 40
    hud_h = 50
    add_rect(slide, 0, hud_y, 1920, hud_h,
             fill_color=C_PANEL, border_color=C_BORDER)
    add_dimension_label(slide, 5, hud_y - 15, 500, 20,
                        f"HUD 바: (0,0) 1920x{hud_h}px, 채우기=#121A2A, Anchor=Top-Stretch")

    # 웨이브 카운터 (좌측)
    add_text(slide, 30, hud_y + 8, 200, 35,
             "웨이브: 3/5", font_size=18, color=C_TEXT_MAIN, bold=True)
    add_dimension_label(slide, 30, hud_y + hud_h + 2, 400, 20,
                        "웨이브: (30,8) 200x35px, 18pt 굵게, #D8E4FF")

    # Bit 카운터 (좌측 중앙)
    add_text(slide, 300, hud_y + 8, 200, 35,
             "Bit: 85", font_size=18, color=C_NEON_GREEN, bold=True)
    add_dimension_label(slide, 300, hud_y + hud_h + 2, 300, 20,
                        "Bit: (300,8) 200x35px, 18pt 굵게, #2BFF88")

    # 기지 HP 바 (우측 중앙~우측)
    hp_label_x = 900
    add_text(slide, hp_label_x, hud_y + 8, 120, 35,
             "기지 HP:", font_size=14, color=C_TEXT_MAIN)
    hp_bar_x, hp_bar_w, hp_bar_h = 1020, 400, 24
    hp_bar_y = hud_y + 13
    # HP 바 배경
    add_rect(slide, hp_bar_x, hp_bar_y, hp_bar_w, hp_bar_h,
             fill_color=RGBColor(0x33, 0x11, 0x11), border_color=C_BORDER)
    # HP 바 채우기 (80%)
    add_rect(slide, hp_bar_x + 2, hp_bar_y + 2, int(hp_bar_w * 0.8) - 4, hp_bar_h - 4,
             fill_color=C_HP_GREEN)
    # HP 바 위 텍스트
    add_text(slide, hp_bar_x, hp_bar_y + 1, hp_bar_w, hp_bar_h,
             "16 / 20", font_size=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_dimension_label(slide, hp_bar_x, hud_y + hud_h + 2, 500, 20,
                        f"HP 바: ({hp_bar_x},13) {hp_bar_w}x{hp_bar_h}px, HP 비율에 비례하여 채움")

    # HP 경고 안내
    add_text(slide, 1450, hud_y + 8, 450, 35,
             "< HP < 30%: 화면 가장자리 붉게 >", font_size=10, color=C_RED)

    # ── 게임 맵 영역 ──
    map_y = 95
    map_h = 850
    add_rect(slide, 0, map_y, 1920, map_h,
             fill_color=RGBColor(0x08, 0x0C, 0x16), border_color=C_BORDER)
    add_dimension_label(slide, 5, map_y + 2, 600, 20,
                        f"게임 맵: (0,{map_y - 40}) 1920x{map_h}px, Stretch-All, 격자 오버레이")

    # 격자선 (샘플)
    grid_size = 80
    grid_start_x = 200
    grid_start_y = map_y + 80
    grid_cols = 18
    grid_rows = 9
    # 개념을 보여주기 위해 몇 개의 격자선 그리기
    for i in range(min(grid_cols + 1, 6)):
        gx = grid_start_x + i * grid_size
        add_rect(slide, gx, grid_start_y, 1, grid_rows * grid_size,
                 fill_color=RGBColor(0x1A, 0x24, 0x3A))
    for j in range(min(grid_rows + 1, 5)):
        gy = grid_start_y + j * grid_size
        add_rect(slide, grid_start_x, gy, 5 * grid_size, 1,
                 fill_color=RGBColor(0x1A, 0x24, 0x3A))
    add_dimension_label(slide, grid_start_x, grid_start_y - 16, 400, 20,
                        f"격자: {grid_size}x{grid_size}px 셀, #1A243A 선, 1px")

    # 격자 위의 샘플 타워
    tower_gx, tower_gy = grid_start_x + grid_size, grid_start_y + grid_size
    add_rounded_rect(slide, tower_gx + 5, tower_gy + 5, grid_size - 10, grid_size - 10,
                     fill_color=C_PANEL, border_color=C_NEON_GREEN, border_width=2)
    add_text(slide, tower_gx + 5, tower_gy + 20, grid_size - 10, 40,
             "Arrow\nLv1", font_size=8, color=C_NEON_GREEN, align=PP_ALIGN.CENTER)

    # 배치 가능 하이라이트 샘플
    empty_gx = grid_start_x + 2 * grid_size
    empty_gy = grid_start_y + grid_size
    add_rect(slide, empty_gx + 2, empty_gy + 2, grid_size - 4, grid_size - 4,
             fill_color=RGBColor(0x15, 0x30, 0x20), border_color=C_NEON_GREEN, border_width=1)
    add_text(slide, empty_gx + 5, empty_gy + 25, grid_size - 10, 30,
             "빈칸", font_size=7, color=C_NEON_GREEN, align=PP_ALIGN.CENTER)

    # 경로 표시
    add_text(slide, 700, map_y + 350, 500, 80,
             "[ Node 경로 ]\n(Node가 이 경로를 따라 기지로 이동)",
             font_size=14, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

    # 기지 아이콘
    add_rounded_rect(slide, 1600, map_y + 380, 80, 80,
                     fill_color=C_PANEL, border_color=C_NEON_BLUE, border_width=2)
    add_text(slide, 1600, map_y + 400, 80, 40,
             "기지", font_size=10, color=C_NEON_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # 스폰 지점
    add_rounded_rect(slide, 100, map_y + 380, 80, 80,
                     fill_color=C_PANEL, border_color=C_RED, border_width=2)
    add_text(slide, 100, map_y + 400, 80, 40,
             "스폰", font_size=9, color=C_RED, bold=True, align=PP_ALIGN.CENTER)

    # ── 하단 바 (타워 인벤토리 + 컨트롤) ──
    bot_y = 950
    bot_h = 95
    add_rect(slide, 0, bot_y, 1920, bot_h,
             fill_color=C_PANEL, border_color=C_BORDER)
    add_dimension_label(slide, 5, bot_y - 15, 600, 20,
                        f"하단 바: (0,{bot_y - 40}) 1920x{bot_h}px, Anchor=Bottom-Stretch")

    # 타워 인벤토리 슬롯
    slot_size = 70
    slot_gap = 10
    slot_start_x = 20
    slot_y = bot_y + 12
    tower_labels = ["Arrow", "Arrow", "Cannon"]
    for i, label in enumerate(tower_labels):
        sx = slot_start_x + i * (slot_size + slot_gap)
        add_rounded_rect(slide, sx, slot_y, slot_size, slot_size,
                         fill_color=C_PANEL_LIGHT, border_color=C_NEON_GREEN, border_width=1)
        add_text(slide, sx, slot_y + 20, slot_size, 30,
                 label, font_size=9, color=C_NEON_GREEN, align=PP_ALIGN.CENTER)
    add_dimension_label(slide, slot_start_x, bot_y + bot_h + 2, 500, 20,
                        f"인벤토리 슬롯: 각 {slot_size}x{slot_size}px, 간격={slot_gap}px, Anchor=Bottom-Left")

    # 구매 버튼
    buy_x = slot_start_x + 3 * (slot_size + slot_gap)
    add_rounded_rect(slide, buy_x, slot_y, slot_size, slot_size,
                     fill_color=C_PANEL_LIGHT, border_color=C_YELLOW, border_width=2)
    add_text(slide, buy_x, slot_y + 15, slot_size, 40,
             "+ 구매", font_size=12, color=C_YELLOW, bold=True, align=PP_ALIGN.CENTER)
    add_dimension_label(slide, buy_x, bot_y + bot_h + 2, 300, 20,
                        f"+구매: ({buy_x},{bot_y - 40 + 12}) {slot_size}x{slot_size}px, #FFD84D")

    # 속도 + 일시정지 컨트롤 (우측)
    ctrl_start_x = 1500
    speed_btns = [("x1", True), ("x2", False), ("x3", False)]
    for i, (label, is_active) in enumerate(speed_btns):
        bx = ctrl_start_x + i * 70
        border = C_NEON_GREEN if is_active else C_BORDER
        fill = RGBColor(0x15, 0x30, 0x20) if is_active else C_PANEL_LIGHT
        add_rounded_rect(slide, bx, slot_y, 60, 35,
                         fill_color=fill, border_color=border)
        add_text(slide, bx, slot_y + 6, 60, 25,
                 label, font_size=12, color=border, bold=is_active, align=PP_ALIGN.CENTER)

    # 일시정지 버튼
    pause_x = ctrl_start_x + 3 * 70
    add_rounded_rect(slide, pause_x, slot_y, 60, 35,
                     fill_color=C_PANEL_LIGHT, border_color=C_ORANGE)
    add_text(slide, pause_x, slot_y + 6, 60, 25,
             "| |", font_size=14, color=C_ORANGE, bold=True, align=PP_ALIGN.CENTER)

    add_dimension_label(slide, ctrl_start_x, slot_y + 40, 400, 20,
                        "속도 버튼: 60x35px, 간격=10px | 일시정지: 60x35px, Anchor=Bottom-Right")

    # 드래그 안내
    add_text(slide, 500, slot_y + 5, 500, 60,
             "인벤토리에서 맵 격자로 드래그하여 배치.\n같은 타워 위에 드래그하면 합성 (레벨 업).",
             font_size=10, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

    # 명세 노트
    notes_data = [
        ("명세 노트 - 인게임", 11, C_ORANGE, True),
        ("HUD 바: 고정 50px, 항상 표시", 9, C_TEXT_DIM, False),
        ("맵: HUD와 하단 바 사이 나머지 공간 채움", 9, C_TEXT_DIM, False),
        ("격자 셀: 80x80px (조정 가능), 내부에 타워 스프라이트", 9, C_TEXT_DIM, False),
        ("인벤토리: 타워 6개 초과 시 가로 스크롤", 9, C_TEXT_DIM, False),
        ("드래그 피드백: 초록=배치 가능, 빨강=불가, 노랑=합성", 9, C_TEXT_DIM, False),
        ("웨이브 클리어: 0.2초 전체화면 플래시, 중앙에 '웨이브 클리어' 텍스트", 9, C_TEXT_DIM, False),
    ]
    add_multiline_text(slide, 1400, map_y + 10, 500, 140, notes_data)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 4: 런 종료 패널 (패배 + 클리어)
# ═══════════════════════════════════════════════════════════════════════════════
def build_slide_run_end(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DARK)

    add_text(slide, 20, 10, 600, 30, "슬라이드 4: 런 종료 패널 (패배 / 스테이지 클리어)",
             font_size=10, color=C_ORANGE, bold=True)

    # ── 반투명 오버레이 ──
    add_rect(slide, 0, 40, 1920, 1040,
             fill_color=RGBColor(0x05, 0x08, 0x12), border_color=C_BORDER)
    add_dimension_label(slide, 5, 42, 600, 20,
                        "오버레이: 전체 화면, #050812 80% 불투명도 (반투명)")

    # ══ 좌측: 패배 ══
    panel_w, panel_h = 800, 700
    defeat_x, defeat_y = 60, 180

    add_rect(slide, defeat_x, defeat_y, panel_w, panel_h,
             fill_color=C_PANEL, border_color=C_RED, border_width=2)
    add_dimension_label(slide, defeat_x, defeat_y - 16, 500, 20,
                        f"패배 패널: ({defeat_x},{defeat_y - 40}) {panel_w}x{panel_h}px, 화면 중앙")

    # 패배 제목
    add_text(slide, defeat_x, defeat_y + 20, panel_w, 60,
             "패배", font_size=36, color=C_RED, bold=True, align=PP_ALIGN.CENTER)
    add_dimension_label(slide, defeat_x + panel_w + 5, defeat_y + 20, 300, 20,
                        "제목: 36pt 굵게, #FF4D5A, 상단여백=20px")

    # 통계
    stats_y = defeat_y + 100
    defeat_stats = [
        ("도달 웨이브:", "4 / 5"),
        ("처치 Node:", "47"),
        ("획득 Bit:", "+185"),
        ("획득 타워:", "[Arrow] [Cannon]"),
    ]
    for i, (label, value) in enumerate(defeat_stats):
        row_y = stats_y + i * 50
        add_text(slide, defeat_x + 80, row_y, 300, 40,
                 label, font_size=16, color=C_TEXT_DIM)
        add_text(slide, defeat_x + 400, row_y, 320, 40,
                 value, font_size=16, color=C_TEXT_MAIN, bold=True)
    add_dimension_label(slide, defeat_x + panel_w + 5, stats_y, 300, 20,
                        "통계: 16pt, 행 간격=50px, 좌측=#AFC3E8, 우측=#D8E4FF 굵게")

    # 패배 버튼
    btn_y = defeat_y + panel_h - 120
    add_rounded_rect(slide, defeat_x + 80, btn_y, 280, 55,
                     fill_color=C_PANEL_LIGHT, border_color=C_NEON_BLUE, border_width=2)
    add_text(slide, defeat_x + 80, btn_y + 12, 280, 35,
             "[ Hub (스킬 트리) ]", font_size=14, color=C_NEON_BLUE,
             bold=True, align=PP_ALIGN.CENTER)

    add_rounded_rect(slide, defeat_x + 420, btn_y, 280, 55,
                     fill_color=C_PANEL_LIGHT, border_color=C_NEON_GREEN, border_width=2)
    add_text(slide, defeat_x + 420, btn_y + 12, 280, 35,
             "[ 즉시 재도전 ]", font_size=14, color=C_NEON_GREEN,
             bold=True, align=PP_ALIGN.CENTER)
    add_dimension_label(slide, defeat_x + 80, btn_y + 60, 500, 20,
                        "버튼: 각 280x55px, 간격=60px, 하단여백=65px")

    # ══ 우측: 스테이지 클리어 ══
    clear_x = 1020
    clear_y = 180

    add_rect(slide, clear_x, clear_y, panel_w, panel_h,
             fill_color=C_PANEL, border_color=C_NEON_GREEN, border_width=2)
    add_dimension_label(slide, clear_x, clear_y - 16, 500, 20,
                        f"클리어 패널: 동일 크기 {panel_w}x{panel_h}px, #2BFF88 테두리")

    # 스테이지 클리어 제목
    add_text(slide, clear_x, clear_y + 20, panel_w, 60,
             "스테이지 클리어!", font_size=36, color=C_NEON_GREEN, bold=True, align=PP_ALIGN.CENTER)

    # 스테이지 이름
    add_text(slide, clear_x, clear_y + 75, panel_w, 30,
             "스테이지 1: 데이터 스트림", font_size=14, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

    # 통계
    clear_stats_y = clear_y + 130
    clear_stats = [
        ("획득 Bit:", "+230"),
        ("획득 Core:", "+2  (신규!)"),
        ("획득 타워:", "[Arrow] [Arrow] [Ice]"),
    ]
    for i, (label, value) in enumerate(clear_stats):
        row_y = clear_stats_y + i * 50
        val_color = C_NEON_PURPLE if "Core" in label else C_TEXT_MAIN
        add_text(slide, clear_x + 80, row_y, 300, 40,
                 label, font_size=16, color=C_TEXT_DIM)
        add_text(slide, clear_x + 400, row_y, 320, 40,
                 value, font_size=16, color=val_color, bold=True)

    # 새 스테이지 해금 알림
    unlock_y = clear_stats_y + 180
    add_rect(slide, clear_x + 60, unlock_y, panel_w - 120, 80,
             fill_color=RGBColor(0x15, 0x20, 0x35), border_color=C_YELLOW, border_width=1)
    add_text(slide, clear_x + 60, unlock_y + 8, panel_w - 120, 30,
             "새 스테이지 해금!", font_size=14, color=C_YELLOW, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, clear_x + 60, unlock_y + 38, panel_w - 120, 30,
             "스테이지 2: 메모리 블록", font_size=16, color=C_TEXT_MAIN, align=PP_ALIGN.CENTER)

    # 클리어 버튼
    btn_y = clear_y + panel_h - 120
    add_rounded_rect(slide, clear_x + 80, btn_y, 280, 55,
                     fill_color=C_PANEL_LIGHT, border_color=C_NEON_BLUE, border_width=2)
    add_text(slide, clear_x + 80, btn_y + 12, 280, 35,
             "[ Hub ]", font_size=14, color=C_NEON_BLUE,
             bold=True, align=PP_ALIGN.CENTER)

    add_rounded_rect(slide, clear_x + 420, btn_y, 280, 55,
                     fill_color=C_PANEL_LIGHT, border_color=C_NEON_GREEN, border_width=2)
    add_text(slide, clear_x + 420, btn_y + 12, 280, 35,
             "[ 다음 스테이지 ]", font_size=14, color=C_NEON_GREEN,
             bold=True, align=PP_ALIGN.CENTER)

    # 명세 노트
    notes_data = [
        ("명세 노트 - 런 종료 패널", 11, C_ORANGE, True),
        ("", 6, C_TEXT_DIM, False),
        ("두 패널 모두: 화면 중앙 (한 번에 하나만 표시)", 9, C_TEXT_DIM, False),
        ("패널 크기: 800x700px, Anchor=Center-Center", 9, C_TEXT_DIM, False),
        ("오버레이: 전체 화면 검정 80% 불투명도", 9, C_TEXT_DIM, False),
        ("애니메이션: 패널이 하단에서 0.3초간 슬라이드 업", 9, C_TEXT_DIM, False),
        ("Core의 '신규!' 라벨: 깜빡이는 애니메이션", 9, C_TEXT_DIM, False),
        ("통계의 타워 아이콘: 50x50px 인라인 스프라이트", 9, C_TEXT_DIM, False),
        ("패배=빨간 테마 | 클리어=초록 테마", 9, C_TEXT_DIM, False),
    ]
    add_multiline_text(slide, 60, 910, 500, 150, notes_data)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 5: 팝업 및 툴팁
# ═══════════════════════════════════════════════════════════════════════════════
def build_slide_popups(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DARK)

    add_text(slide, 20, 10, 500, 30, "슬라이드 5: 팝업 및 툴팁",
             font_size=10, color=C_ORANGE, bold=True)

    # ══ 1. 스킬 노드 구매 팝업 (좌측 상단) ══
    px1, py1 = 40, 60
    pw, ph = 380, 340
    add_rect(slide, px1, py1, pw, ph,
             fill_color=C_PANEL, border_color=C_NEON_GREEN, border_width=2)
    add_dimension_label(slide, px1, py1 - 15, 400, 20,
                        f"스킬 구매 팝업: {pw}x{ph}px, 열릴 때 화면 중앙")

    # 제목
    add_text(slide, px1 + 20, py1 + 15, pw - 40, 30,
             "공격력  Lv3 -> Lv4", font_size=16, color=C_NEON_GREEN, bold=True)
    # 구분선
    add_rect(slide, px1 + 20, py1 + 55, pw - 40, 1, fill_color=C_BORDER)
    # 효과
    add_text(slide, px1 + 20, py1 + 65, pw - 40, 25,
             "효과: 타워 피해량 +10%", font_size=13, color=C_TEXT_MAIN)
    # 변경 전/후
    add_text(slide, px1 + 20, py1 + 100, pw - 40, 25,
             "현재: x1.3  ->  x1.4", font_size=13, color=C_TEXT_DIM)
    # 비용
    add_text(slide, px1 + 20, py1 + 145, pw - 40, 25,
             "비용: 110 Bit", font_size=15, color=C_YELLOW, bold=True)
    # 구매 버튼
    add_rounded_rect(slide, px1 + 60, py1 + 200, 260, 50,
                     fill_color=RGBColor(0x15, 0x30, 0x20), border_color=C_NEON_GREEN, border_width=2)
    add_text(slide, px1 + 60, py1 + 210, 260, 35,
             "[ 구매 ]", font_size=16, color=C_NEON_GREEN,
             bold=True, align=PP_ALIGN.CENTER)
    # 닫기 X
    add_text(slide, px1 + pw - 40, py1 + 10, 30, 25,
             "X", font_size=14, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

    # 구매 팝업 명세 노트
    popup_notes = [
        ("스킬 구매 팝업", 9, C_ORANGE, True),
        ("크기: 380x340px", 8, C_TEXT_DIM, False),
        ("Anchor: Center-Center (오버레이)", 8, C_TEXT_DIM, False),
        ("제목: 16pt 굵게 #2BFF88", 8, C_TEXT_DIM, False),
        ("효과/변경 전후: 13pt #D8E4FF/#AFC3E8", 8, C_TEXT_DIM, False),
        ("비용: 15pt 굵게 #FFD84D", 8, C_TEXT_DIM, False),
        ("구매 버튼: 260x50px, 하단 중앙", 8, C_TEXT_DIM, False),
        ("닫기 X: 30x25px, 우측 상단", 8, C_TEXT_DIM, False),
        ("자금 부족 시: 비용 텍스트가 #FF4D5A로 변경", 8, C_TEXT_DIM, False),
        ("최대 레벨: 버튼 비활성화, '최대' 텍스트", 8, C_TEXT_DIM, False),
    ]
    add_multiline_text(slide, px1, py1 + ph + 10, pw, 180, popup_notes)

    # ══ 2. 타워 정보 툴팁 (상단 중앙) ══
    tx1, ty1 = 500, 60
    tw, th = 350, 300
    add_rect(slide, tx1, ty1, tw, th,
             fill_color=C_PANEL, border_color=C_NEON_BLUE, border_width=2)
    add_dimension_label(slide, tx1, ty1 - 15, 400, 20,
                        f"타워 툴팁: {tw}x{th}px, 타워 호버/클릭 시 인접 표시")

    # 타워 이름 + 레벨
    add_text(slide, tx1 + 20, ty1 + 15, tw - 40, 30,
             "Arrow 타워  Lv2", font_size=16, color=C_NEON_BLUE, bold=True)
    add_rect(slide, tx1 + 20, ty1 + 50, tw - 40, 1, fill_color=C_BORDER)

    # 스탯
    tower_stats = [
        ("피해량:", "12"),
        ("공격 속도:", "0.7초"),
        ("사거리:", "3.2"),
        ("DPS:", "17.1"),
    ]
    for i, (label, val) in enumerate(tower_stats):
        row_y = ty1 + 60 + i * 32
        add_text(slide, tx1 + 20, row_y, 160, 25,
                 label, font_size=12, color=C_TEXT_DIM)
        add_text(slide, tx1 + 200, row_y, 120, 25,
                 val, font_size=12, color=C_TEXT_MAIN, bold=True)

    # 특수 효과
    add_text(slide, tx1 + 20, ty1 + 195, tw - 40, 25,
             "특수: -", font_size=11, color=C_TEXT_DIM)
    # 판매 정보
    add_text(slide, tx1 + 20, ty1 + 235, tw - 40, 25,
             "판매: 17 Bit (70%)", font_size=11, color=C_ORANGE)

    # 툴팁 노트
    tooltip_notes = [
        ("타워 툴팁", 9, C_ORANGE, True),
        ("크기: 350x300px", 8, C_TEXT_DIM, False),
        ("위치: 타워 인접, 화면 가장자리 회피", 8, C_TEXT_DIM, False),
        ("스탯: 12pt, 라벨=#AFC3E8, 값=#D8E4FF 굵게", 8, C_TEXT_DIM, False),
        ("판매 가격: 11pt #FF9A3D", 8, C_TEXT_DIM, False),
        ("표시 조건: 타워 클릭 (모바일/데스크톱)", 8, C_TEXT_DIM, False),
    ]
    add_multiline_text(slide, tx1, ty1 + th + 10, tw, 120, tooltip_notes)

    # ══ 3. 설정 팝업 (우측 상단) ══
    sx1, sy1 = 950, 60
    sw, sh = 400, 300
    add_rect(slide, sx1, sy1, sw, sh,
             fill_color=C_PANEL, border_color=C_BORDER, border_width=2)
    add_dimension_label(slide, sx1, sy1 - 15, 400, 20,
                        f"설정 팝업: {sw}x{sh}px, 화면 중앙")

    add_text(slide, sx1, sy1 + 15, sw, 30,
             "설정", font_size=20, color=C_TEXT_MAIN, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, sx1 + 30, sy1 + 55, sw - 60, 1, fill_color=C_BORDER)

    # BGM 슬라이더
    add_text(slide, sx1 + 30, sy1 + 75, 100, 25,
             "BGM", font_size=14, color=C_TEXT_MAIN)
    slider_x = sx1 + 130
    add_rect(slide, slider_x, sy1 + 82, 220, 8,
             fill_color=C_PANEL_LIGHT, border_color=C_BORDER)
    add_rect(slide, slider_x, sy1 + 82, 154, 8,
             fill_color=C_NEON_BLUE)
    add_rounded_rect(slide, slider_x + 150, sy1 + 76, 16, 20,
                     fill_color=C_WHITE, border_color=C_NEON_BLUE)
    add_text(slide, sx1 + sw - 60, sy1 + 75, 40, 25,
             "70%", font_size=11, color=C_TEXT_DIM)

    # SFX 슬라이더
    add_text(slide, sx1 + 30, sy1 + 125, 100, 25,
             "SFX", font_size=14, color=C_TEXT_MAIN)
    add_rect(slide, slider_x, sy1 + 132, 220, 8,
             fill_color=C_PANEL_LIGHT, border_color=C_BORDER)
    add_rect(slide, slider_x, sy1 + 132, 198, 8,
             fill_color=C_NEON_BLUE)
    add_rounded_rect(slide, slider_x + 194, sy1 + 126, 16, 20,
                     fill_color=C_WHITE, border_color=C_NEON_BLUE)
    add_text(slide, sx1 + sw - 60, sy1 + 125, 40, 25,
             "90%", font_size=11, color=C_TEXT_DIM)

    # 닫기 버튼
    add_rounded_rect(slide, sx1 + 120, sy1 + sh - 70, 160, 45,
                     fill_color=C_PANEL_LIGHT, border_color=C_BORDER, border_width=1)
    add_text(slide, sx1 + 120, sy1 + sh - 60, 160, 35,
             "[ 닫기 ]", font_size=14, color=C_TEXT_MAIN, align=PP_ALIGN.CENTER)

    # 닫기 X
    add_text(slide, sx1 + sw - 40, sy1 + 10, 30, 25,
             "X", font_size=14, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

    # 설정 노트
    settings_notes = [
        ("설정 팝업", 9, C_ORANGE, True),
        ("크기: 400x300px, Center-Center", 8, C_TEXT_DIM, False),
        ("슬라이더 트랙: 220x8px, #1A243A 배경", 8, C_TEXT_DIM, False),
        ("슬라이더 채움: #37B6FF", 8, C_TEXT_DIM, False),
        ("슬라이더 핸들: 16x20px, 흰색", 8, C_TEXT_DIM, False),
        ("닫기 버튼: 160x45px, 하단 중앙", 8, C_TEXT_DIM, False),
    ]
    add_multiline_text(slide, sx1, sy1 + sh + 10, sw, 120, settings_notes)

    # ══ 4. 방치 Bit 수령 팝업 (우측) ══
    ix1, iy1 = 1450, 60
    iw, ih = 420, 250
    add_rect(slide, ix1, iy1, iw, ih,
             fill_color=C_PANEL, border_color=C_YELLOW, border_width=2)
    add_dimension_label(slide, ix1, iy1 - 15, 400, 20,
                        f"방치 수령 팝업: {iw}x{ih}px, 화면 중앙 또는 상단 중앙")

    add_text(slide, ix1, iy1 + 20, iw, 30,
             "방치 보상", font_size=20, color=C_YELLOW, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, ix1 + 30, iy1 + 60, iw - 60, 1, fill_color=C_BORDER)

    add_text(slide, ix1, iy1 + 80, iw, 40,
             "+150 Bit", font_size=28, color=C_NEON_GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, ix1, iy1 + 120, iw, 25,
             "부재 시간 1시간 30분", font_size=12, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

    # 수령 버튼
    add_rounded_rect(slide, ix1 + 100, iy1 + 170, 220, 50,
                     fill_color=RGBColor(0x30, 0x28, 0x10), border_color=C_YELLOW, border_width=2)
    add_text(slide, ix1 + 100, iy1 + 180, 220, 35,
             "[ 수령 ]", font_size=16, color=C_YELLOW,
             bold=True, align=PP_ALIGN.CENTER)

    # 방치 노트
    idle_notes = [
        ("방치 수령 팝업", 9, C_ORANGE, True),
        ("크기: 420x250px", 8, C_TEXT_DIM, False),
        ("표시: Hub 진입 시 방치 Bit > 0이면", 8, C_TEXT_DIM, False),
        ("Bit 금액: 28pt 굵게 #2BFF88, 카운트업 애니메이션", 8, C_TEXT_DIM, False),
        ("수령 버튼: 220x50px, #FFD84D 테두리", 8, C_TEXT_DIM, False),
        ("자동 닫힘: 수령 후, +Bit 팝업이 위로 떠오름", 8, C_TEXT_DIM, False),
    ]
    add_multiline_text(slide, ix1, iy1 + ih + 10, iw, 120, idle_notes)

    # ══ 타워 구매 팝업 (하단 영역) ══
    bx1, by1 = 500, 680
    bw, bh = 500, 340
    add_rect(slide, bx1, by1, bw, bh,
             fill_color=C_PANEL, border_color=C_YELLOW, border_width=2)
    add_dimension_label(slide, bx1, by1 - 15, 500, 20,
                        f"타워 구매 패널: {bw}x{bh}px, 인게임 +구매 버튼 위에 표시")

    add_text(slide, bx1, by1 + 10, bw, 30,
             "타워 구매", font_size=18, color=C_YELLOW, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, bx1 + 20, by1 + 45, bw - 40, 1, fill_color=C_BORDER)

    # 타워 선택지
    tower_options = [
        ("Arrow 타워", "30 Bit", C_NEON_GREEN),
        ("Cannon 타워", "60 Bit", C_NEON_GREEN),
        ("Ice 타워", "45 Bit", C_NEON_GREEN),
    ]
    for i, (name, cost, color) in enumerate(tower_options):
        oy = by1 + 55 + i * 65
        add_rounded_rect(slide, bx1 + 30, oy, bw - 60, 55,
                         fill_color=C_PANEL_LIGHT, border_color=C_BORDER)
        add_text(slide, bx1 + 50, oy + 12, 200, 30,
                 name, font_size=14, color=C_TEXT_MAIN, bold=True)
        add_rounded_rect(slide, bx1 + bw - 170, oy + 8, 120, 38,
                         fill_color=RGBColor(0x15, 0x30, 0x20), border_color=color)
        add_text(slide, bx1 + bw - 170, oy + 14, 120, 30,
                 cost, font_size=12, color=color, bold=True, align=PP_ALIGN.CENTER)

    # 구매 패널 노트
    buy_notes = [
        ("타워 구매 패널", 9, C_ORANGE, True),
        ("크기: 500x340px (타워 추가 시 확장)", 8, C_TEXT_DIM, False),
        ("위치: +구매 버튼 위, Anchor=Bottom-Left", 8, C_TEXT_DIM, False),
        ("스킬 트리에서 해금된 타워만 표시", 8, C_TEXT_DIM, False),
        ("각 행: 높이 55px, 타워 아이콘 + 이름 + 비용 버튼", 8, C_TEXT_DIM, False),
        ("잠긴 타워: 회색 처리 + 잠금 아이콘", 8, C_TEXT_DIM, False),
    ]
    add_multiline_text(slide, bx1 + bw + 20, by1, 400, 120, buy_notes)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 6: 색상 팔레트 참조
# ═══════════════════════════════════════════════════════════════════════════════
def build_slide_palette(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_DARK)

    add_text(slide, 20, 10, 700, 30,
             "슬라이드 6: 색상 팔레트 및 타이포그래피 참조 (네온 회로판 테마)",
             font_size=10, color=C_ORANGE, bold=True)

    # ── 색상 팔레트 섹션 ──
    add_text(slide, 40, 50, 400, 30,
             "색상 팔레트", font_size=18, color=C_TEXT_MAIN, bold=True)

    colors = [
        ("#0B0F1A", "어두운 배경", "메인 배경, 게임 영역", C_BG_DARK, C_WHITE),
        ("#121A2A", "패널", "패널/카드 배경", C_PANEL, C_WHITE),
        ("#1A243A", "밝은 패널", "호버 상태, 밝은 패널", C_PANEL_LIGHT, C_WHITE),
        ("#5B6B8A", "테두리", "기본 테두리, 구분선", C_BORDER, C_WHITE),
        ("#D8E4FF", "기본 텍스트", "기본 텍스트, 제목", C_TEXT_MAIN, C_BG_DARK),
        ("#AFC3E8", "보조 텍스트", "보조 텍스트, 라벨", C_TEXT_DIM, C_BG_DARK),
        ("#2BFF88", "네온 초록", "기본 강조, 활성, 성공, HP", C_NEON_GREEN, C_BG_DARK),
        ("#37B6FF", "네온 파랑", "보조 강조, 정보, 링크", C_NEON_BLUE, C_BG_DARK),
        ("#B86CFF", "네온 보라", "Core 자원, 프리미엄", C_NEON_PURPLE, C_BG_DARK),
        ("#FF9A3D", "주황", "경고, 하이라이트, 판매", C_ORANGE, C_BG_DARK),
        ("#FF4D5A", "빨강", "경보, 위험, 패배, HP 부족", C_RED, C_BG_DARK),
        ("#FFD84D", "노랑/금색", "희귀, 방치 보상, 구매", C_YELLOW, C_BG_DARK),
    ]

    chip_w, chip_h = 180, 55
    label_w = 280
    for i, (hex_code, name, usage, fill, text_c) in enumerate(colors):
        row = i % 6
        col = i // 6
        x = 40 + col * 900
        y = 90 + row * 65
        # 색상 칩
        add_rect(slide, x, y, chip_w, chip_h, fill_color=fill, border_color=C_BORDER)
        add_text(slide, x + 10, y + 8, chip_w - 20, 20,
                 hex_code, font_size=10, color=text_c, bold=True)
        add_text(slide, x + 10, y + 28, chip_w - 20, 20,
                 name, font_size=9, color=text_c)
        # 용도 설명
        add_text(slide, x + chip_w + 10, y + 12, label_w, 35,
                 usage, font_size=9, color=C_TEXT_DIM)

    # ── 타이포그래피 섹션 ──
    typo_y = 500
    add_text(slide, 40, typo_y, 400, 30,
             "타이포그래피", font_size=18, color=C_TEXT_MAIN, bold=True)

    font_specs = [
        ("48pt 굵게", "게임 제목 (타이틀 화면)", 48, True),
        ("36pt 굵게", "주요 제목 (패배, 스테이지 클리어)", 22, True),
        ("20pt 굵게", "팝업 제목 (설정, 방치 보상)", 20, True),
        ("18pt 굵게", "HUD 값 (웨이브, Bit, HP)", 18, True),
        ("16pt 굵게", "소제목, 버튼 텍스트, 스탯 값", 16, True),
        ("14pt", "본문 텍스트, 스탯 라벨, 설명", 14, False),
        ("12pt", "작은 라벨, 드롭다운 텍스트, 슬라이더 값", 12, False),
        ("10pt", "주석, 툴팁, 보조 정보", 10, False),
        ("8pt", "명세 노트, 치수 라벨 (개발 전용)", 8, False),
    ]

    for i, (size_label, usage, pt, bold) in enumerate(font_specs):
        row_y = typo_y + 35 + i * 32
        add_text(slide, 60, row_y, 160, 28,
                 size_label, font_size=min(pt, 14), color=C_TEXT_MAIN, bold=bold)
        add_text(slide, 240, row_y, 600, 28,
                 usage, font_size=10, color=C_TEXT_DIM)

    # ── UI 요소 크기 요약 ──
    sizes_x = 960
    sizes_y = 500
    add_text(slide, sizes_x, sizes_y, 400, 30,
             "표준 요소 크기", font_size=18, color=C_TEXT_MAIN, bold=True)

    size_specs = [
        "버튼 (대): 300x55px, 테두리=2px, 라운드=8px",
        "버튼 (중): 160x45px, 테두리=1px, 라운드=6px",
        "버튼 (소): 100x38px, 테두리=1px, 라운드=4px",
        "속도 버튼: 60x35px",
        "타워 슬롯: 70x70px, 간격=10px",
        "스킬 트리 노드: 60x60px, 중심 간격=80px",
        "격자 셀: 80x80px",
        "HP 바: 400x24px",
        "슬라이더 트랙: 220x8px, 핸들=16x20px",
        "상단/하단 바: 전체 너비 x 50~55px",
        "팝업 안쪽 여백: 사방 20~30px",
        "패널 테두리: 1~2px",
    ]

    for i, spec in enumerate(size_specs):
        row_y = sizes_y + 35 + i * 26
        add_text(slide, sizes_x + 20, row_y, 500, 24,
                 spec, font_size=9, color=C_TEXT_DIM)

    # ── Anchor 참조 ──
    anchor_y = 870
    add_text(slide, 40, anchor_y, 600, 30,
             "Unity RectTransform Anchor 참조", font_size=14, color=C_TEXT_MAIN, bold=True)
    anchors = [
        "배경/오버레이: Stretch-Stretch (모든 가장자리 = 0)",
        "상단 바: Top-Stretch (left=0, right=0, top=0, height=50)",
        "하단 바: Bottom-Stretch (left=0, right=0, bottom=0, height=55)",
        "중앙 패널 (팝업): Center-Center, Pivot=(0.5, 0.5)",
        "HUD 요소: Top-Left 또는 Top-Right + 오프셋",
        "인벤토리 슬롯: Bottom-Left, 가로 배치",
        "속도/일시정지: Bottom-Right, 가로 배치",
        "스킬 트리 뷰포트: Stretch-All + ScrollRect",
    ]
    for i, anchor in enumerate(anchors):
        add_text(slide, 60, anchor_y + 30 + i * 22, 800, 20,
                 anchor, font_size=9, color=C_TEXT_DIM)


# ═══════════════════════════════════════════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()

    # 슬라이드 크기를 1920x1080 px로 설정
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_title_screen(prs)
    build_slide_hub(prs)
    build_slide_ingame(prs)
    build_slide_run_end(prs)
    build_slide_popups(prs)
    build_slide_palette(prs)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"생성 완료: {OUT_PATH}")
    print(f"슬라이드 수: {len(prs.slides)}")
    print(f"파일 크기: {os.path.getsize(OUT_PATH) / 1024:.1f} KB")


if __name__ == "__main__":
    main()
