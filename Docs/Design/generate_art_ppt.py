#!/usr/bin/env python3
"""Soulspire Art Direction v0.1 PPT Generator"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === Constants ===
BG_COLOR = RGBColor(0x14, 0x14, 0x20)       # 어두운 남보라
TEXT_MAIN = RGBColor(0xE0, 0xDC, 0xD0)       # 밝은 크림
TEXT_ACCENT = RGBColor(0xFF, 0xD8, 0x4D)     # 금색
TEXT_SUB = RGBColor(0xA0, 0x98, 0x90)        # 서브 텍스트
FRAME_COLOR = RGBColor(0x5A, 0x50, 0x70)     # 중간 보라
DARK_BG = RGBColor(0x0A, 0x0A, 0x12)        # 최심부 배경

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_MAIN, bold=False, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_color_rect(slide, left, top, width, height, color_hex, label=None):
    """Add a colored rectangle with optional label"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color_hex)
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = FRAME_COLOR
    shape.line.width = Pt(1)

    if label:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MAIN
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_divider(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.5), top, Inches(12.333), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = FRAME_COLOR
    shape.line.fill.background()


def add_multiline_text(slide, left, top, width, height, lines, default_size=16,
                       default_color=TEXT_MAIN):
    """lines: list of (text, font_size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Arial'
        p.space_after = Pt(4)
    return tf


# ============================================================
# SLIDE 1: 표지
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Decorative frame
frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(1), Inches(0.8), Inches(11.333), Inches(5.9))
frame.fill.solid()
frame.fill.fore_color.rgb = BG_COLOR
frame.line.fill.solid()
frame.line.fill.fore_color.rgb = TEXT_ACCENT
frame.line.width = Pt(3)

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(1.2),
            'SOULSPIRE', 60, TEXT_ACCENT, True, PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10.333), Inches(0.8),
            'Art Direction v0.1', 32, TEXT_MAIN, False, PP_ALIGN.CENTER)

add_divider(slide, Inches(3.8))

add_multiline_text(slide, Inches(2), Inches(4.2), Inches(9.333), Inches(2), [
    ('Dark Fantasy Pixel Art Tower Defense', 24, TEXT_SUB, False),
    ('', 12),
    ('어둠 속에서 빛나는 마법, 파괴의 쾌감을 눈으로 체감하라', 20, TEXT_MAIN, False),
    ('', 12),
    ('2026.02.14  |  TA팀장 (아트디렉터)  |  v0.1 초판', 14, TEXT_SUB, False),
], default_color=TEXT_SUB)

# ============================================================
# SLIDE 2: 아트 비전 요약
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '01  ART VISION', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

add_multiline_text(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(2.5), [
    ('Soulspire는 다크 판타지 세계관의 픽셀 아트 타워 디펜스이다.', 20, TEXT_MAIN, True),
    ('', 8),
    ('어둡고 깊은 톤의 판타지 세계에서, 마법 타워들이 밀려오는 유기체 몬스터를 파괴하는', 16, TEXT_MAIN, False),
    ('시각적 쾌감을 극대화한다. 3/4 쿼터뷰 픽셀 아트로 탑다운 TD의 전략적 가독성과', 16, TEXT_MAIN, False),
    ('디테일한 스프라이트 표현을 동시에 확보하며, 미니멀한 배경 위에서 타워와 몬스터,', 16, TEXT_MAIN, False),
    ('그리고 화려한 마법 이펙트가 시각적 주인공이 되는 구성을 추구한다.', 16, TEXT_MAIN, False),
])

# 핵심 키워드 boxes
keywords = [
    ('다크 판타지', '#9060D0'),
    ('픽셀 아트', '#4080D4'),
    ('3/4 쿼터뷰', '#40D470'),
    ('마법 이펙트', '#FFD84D'),
    ('성장 체감', '#E08030'),
    ('미니멀 배경', '#5A5070'),
]

y = Inches(4.2)
for i, (kw, color) in enumerate(keywords):
    x = Inches(0.8 + i * 2.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Inches(1.8), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = kw
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER

# 톤 매트릭스
add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.5),
            'Tone Matrix', 20, TEXT_ACCENT, True)

tone_items = [
    ('밝음 ←→ 어두움', '70~80% 어두움', '배경과 기본 톤은 어둡다'),
    ('가벼움 ←→ 무거움', '60% 무거움', '시리어스하지만 빠르고 경쾌'),
    ('사실적 ←→ 양식적', '80% 양식적', '픽셀 아트의 양식적 표현'),
    ('귀여움 ←→ 위엄', '65% 위엄', '타워는 위엄, 몬스터는 위협적'),
]

for i, (axis, pos, desc) in enumerate(tone_items):
    x = Inches(0.5 + i * 3.1)
    add_multiline_text(slide, x, Inches(5.8), Inches(2.8), Inches(1.2), [
        (axis, 12, TEXT_ACCENT, True),
        (f'{pos} — {desc}', 11, TEXT_SUB, False),
    ])

# ============================================================
# SLIDE 3: 세계관 전환
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '02  WORLD TRANSITION', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
            '사이버 네온 → 다크 판타지 세계관 전면 전환', 18, TEXT_MAIN, False)

# Comparison table
transitions = [
    ('Node = 데이터 패킷', 'Node = 마계/차원의 유기체 몬스터'),
    ('Tower = 방화벽/보안', 'Tower = 마법 타워/마탑'),
    ('Bit = 데이터 자원', 'Bit = 마나 결정/소울 파편'),
    ('Core = 프리미엄 데이터', 'Core = 고대 룬석/정수'),
    ('회로기판 맵', '황폐한 마법진/고대 유적'),
    ('네온 발광', '마법 발광/룬 광채'),
]

# Table headers
header_y = Inches(1.9)
# Left column header (old)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(1), header_y, Inches(5), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x30, 0x20, 0x40)
shape.line.fill.solid()
shape.line.fill.fore_color.rgb = FRAME_COLOR
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = '기존 (사이버 네온)'
p.font.size = Pt(14)
p.font.color.rgb = TEXT_SUB
p.font.bold = True
p.font.name = 'Arial'
p.alignment = PP_ALIGN.CENTER

# Right column header (new)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(6.2), header_y, Inches(5), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x38)
shape.line.fill.solid()
shape.line.fill.fore_color.rgb = TEXT_ACCENT
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = '신규 (다크 판타지)'
p.font.size = Pt(14)
p.font.color.rgb = TEXT_ACCENT
p.font.bold = True
p.font.name = 'Arial'
p.alignment = PP_ALIGN.CENTER

# Arrow
add_textbox(slide, Inches(5.9), header_y, Inches(0.4), Inches(0.5),
            '→', 20, TEXT_ACCENT, True, PP_ALIGN.CENTER)

for i, (old, new) in enumerate(transitions):
    row_y = Inches(2.5 + i * 0.55)
    # old
    add_textbox(slide, Inches(1.2), row_y, Inches(4.8), Inches(0.5),
                old, 14, TEXT_SUB, False)
    # arrow
    add_textbox(slide, Inches(5.9), row_y, Inches(0.4), Inches(0.5),
                '→', 14, TEXT_ACCENT, True, PP_ALIGN.CENTER)
    # new
    add_textbox(slide, Inches(6.4), row_y, Inches(4.8), Inches(0.5),
                new, 14, TEXT_MAIN, True)

# 핵심 세계관 키워드
add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12), Inches(0.5),
            '핵심 키워드', 18, TEXT_ACCENT, True)

world_kw = [
    '다크 판타지: 신비롭고 위험한 세계',
    '마법 요새: 차원의 균열에서 밀려오는 몬스터를 마법 타워로 방어',
    '고대 유적: 미니멀하지만 신비로운 맵',
    '성장 = 마법 각성: 스킬 트리는 마법 각성 트리',
]
for i, kw in enumerate(world_kw):
    add_textbox(slide, Inches(0.8 + (i % 2) * 6), Inches(6.4 + (i // 2) * 0.4),
                Inches(5.8), Inches(0.4), f'• {kw}', 12, TEXT_MAIN, False)


# ============================================================
# SLIDE 4: 아트 스타일
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '03  ART STYLE', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

# Style spec table
style_items = [
    ('렌더링 스타일', '픽셀 아트 (AI 생성 + 수작업 보정)'),
    ('시점', '3/4 쿼터뷰 (아이소메트릭, ~30-45도)'),
    ('톤', '다크 판타지 (어두운 배경 + 밝은 마법 하이라이트)'),
    ('라이팅', '주변광 없음. 마법 효과와 자체 발광만 존재'),
    ('외곽선', '1~2px 검정 외곽선 (가독성 확보)'),
    ('안티앨리어싱', '수동 AA (제한된 팔레트 내 중간색)'),
]

for i, (label, value) in enumerate(style_items):
    y = Inches(1.3 + i * 0.55)
    add_textbox(slide, Inches(0.8), y, Inches(3), Inches(0.5),
                label, 14, TEXT_ACCENT, True)
    add_textbox(slide, Inches(4), y, Inches(8), Inches(0.5),
                value, 14, TEXT_MAIN, False)

# 3/4 쿼터뷰 설명 box
qv_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.2))
qv_box.fill.solid()
qv_box.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x28)
qv_box.line.fill.solid()
qv_box.line.fill.fore_color.rgb = FRAME_COLOR

add_multiline_text(slide, Inches(1.0), Inches(4.9), Inches(5), Inches(2.0), [
    ('3/4 Quarter View', 16, TEXT_ACCENT, True),
    ('', 6),
    ('• 오브젝트의 정면(앞면)과 윗면이 동시에 보임', 12, TEXT_MAIN, False),
    ('• 높이감 표현 가능 (타워 높이, 몬스터 크기)', 12, TEXT_MAIN, False),
    ('• 탑다운보다 입체감, 사이드뷰보다 전략적 가독성', 12, TEXT_MAIN, False),
    ('• 모든 오브젝트가 동일한 소실점 공유', 12, TEXT_MAIN, False),
])

# Reference list
ref_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(6.8), Inches(4.8), Inches(5.7), Inches(2.2))
ref_box.fill.solid()
ref_box.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x28)
ref_box.line.fill.solid()
ref_box.line.fill.fore_color.rgb = FRAME_COLOR

add_multiline_text(slide, Inches(7.0), Inches(4.9), Inches(5.3), Inches(2.0), [
    ('Visual References', 16, TEXT_ACCENT, True),
    ('', 6),
    ('• Kingdom: Two Crowns — 다크 판타지 + 픽셀 아트', 11, TEXT_MAIN, False),
    ('• Loop Hero — 미니멀 맵 + 성장 루프', 11, TEXT_MAIN, False),
    ('• Vampire Survivors — 이펙트 밀도 스케일링', 11, TEXT_MAIN, False),
    ('• Darkest Dungeon — 다크 판타지 톤, UI', 11, TEXT_MAIN, False),
    ('• Shovel Knight — 픽셀 아트 품질 벤치마크', 11, TEXT_MAIN, False),
])


# ============================================================
# SLIDE 5: 컬러 팔레트
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '04  COLOR PALETTE', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

# --- Background/Environment colors ---
add_textbox(slide, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
            '배경/환경 (어두운 톤)', 16, TEXT_ACCENT, True)

bg_colors = [
    ('#0A0A12', '최심부 배경'),
    ('#141420', '기본 배경'),
    ('#1E1E30', '밝은 배경'),
    ('#2A2A3A', '돌/유적'),
    ('#3A3A50', '밝은 돌'),
]

for i, (hex_c, label) in enumerate(bg_colors):
    x = Inches(0.5 + i * 1.3)
    add_color_rect(slide, x, Inches(1.6), Inches(1.1), Inches(0.7), hex_c)
    add_textbox(slide, x, Inches(2.35), Inches(1.1), Inches(0.5),
                f'{label}\n{hex_c}', 9, TEXT_SUB, False, PP_ALIGN.CENTER)

# --- Magic/Glow colors ---
add_textbox(slide, Inches(0.5), Inches(3.1), Inches(6), Inches(0.4),
            '마법/발광 (하이라이트)', 16, TEXT_ACCENT, True)

magic_colors = [
    ('#E8E4F0', '마법 백색'),
    ('#FFD84D', '골드'),
    ('#D44040', '루비 레드'),
    ('#40D470', '에메랄드'),
    ('#4080D4', '사파이어'),
    ('#9060D0', '마법 퍼플'),
    ('#E08030', '화염 오렌지'),
    ('#E0D040', '번개 옐로우'),
]

for i, (hex_c, label) in enumerate(magic_colors):
    x = Inches(0.5 + i * 1.55)
    add_color_rect(slide, x, Inches(3.55), Inches(1.3), Inches(0.7), hex_c)
    add_textbox(slide, x, Inches(4.3), Inches(1.3), Inches(0.5),
                f'{label}\n{hex_c}', 9, TEXT_SUB, False, PP_ALIGN.CENTER)

# --- Palette principles ---
add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12), Inches(0.4),
            '팔레트 운용 원칙', 16, TEXT_ACCENT, True)

principles = [
    '암부가 기본, 명부는 마법 — 화면 80%는 어두운 톤',
    '채도 제한 — 배경/환경은 저채도, 마법/이펙트만 고채도',
    '스테이지별 색온도 변화 — 초반 차가운 톤, 후반 따뜻한 톤',
    '타워별 즉시 식별 — 시그니처 컬러만으로 종류 구분',
    '몬스터 가독성 — 어두운 배경 위에서 실루엣 즉시 인식',
]

for i, p_text in enumerate(principles):
    add_textbox(slide, Inches(0.8 + (i % 2) * 6), Inches(5.5 + (i // 2) * 0.35),
                Inches(5.8), Inches(0.35), f'• {p_text}', 11, TEXT_MAIN, False)


# ============================================================
# SLIDE 6: 타워 가이드
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '05  TOWER GUIDE (6종)', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

towers = [
    ('T01', '마법 석궁탑', 'Arrow Tower', '#C0A870', '#E8E4F0',
     '석재 기단 + 마법 석궁\n황동빛 금속'),
    ('T02', '화염 마법포', 'Cannon Tower', '#E08030', '#D44040',
     '화산석 포탑\n포구에서 화염 오라'),
    ('T03', '빙결 마탑', 'Ice Tower', '#4080D4', '#80C0E0',
     '얼음 결정 첨탑\n주변에 서리 오라'),
    ('T04', '뇌전 첨탑', 'Lightning Tower', '#E0D040', '#FFFFFF',
     '뾰족한 금속 첨탑\n번개 아크 오라'),
    ('T05', '마력 포탑', 'Laser Tower', '#9060D0', '#D080FF',
     '보라빛 수정구 탑\n마력 빔 집중'),
    ('T06', '공허의 오벨리스크', 'Void Tower', '#302050', '#6040A0',
     '검보라 돌기둥\n주변 공기가 왜곡'),
]

for i, (tid, name, eng, main_c, sub_c, desc) in enumerate(towers):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.1)
    y = Inches(1.2 + row * 3.0)

    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(3.8), Inches(2.7))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x28)
    card.line.fill.solid()
    card.line.fill.fore_color.rgb = FRAME_COLOR

    # Tower ID & name
    add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.5), Inches(0.4),
                f'{tid}  {name}', 16, TEXT_ACCENT, True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.5), Inches(3.5), Inches(0.3),
                eng, 12, TEXT_SUB, False)

    # Main color rect
    add_color_rect(slide, x + Inches(0.15), y + Inches(0.9), Inches(1.2), Inches(0.6), main_c)
    add_textbox(slide, x + Inches(0.15), y + Inches(1.55), Inches(1.2), Inches(0.3),
                f'주색 {main_c}', 9, TEXT_SUB, False, PP_ALIGN.CENTER)

    # Sub color rect
    add_color_rect(slide, x + Inches(1.5), y + Inches(0.9), Inches(1.0), Inches(0.6), sub_c)
    add_textbox(slide, x + Inches(1.5), y + Inches(1.55), Inches(1.0), Inches(0.3),
                f'보조 {sub_c}', 9, TEXT_SUB, False, PP_ALIGN.CENTER)

    # Description
    add_textbox(slide, x + Inches(0.15), y + Inches(1.9), Inches(3.5), Inches(0.7),
                desc, 11, TEXT_MAIN, False)


# ============================================================
# SLIDE 7: 몬스터 가이드
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '06  MONSTER GUIDE (9종)', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

monsters = [
    ('N01', '소형 마물', 'Bit Node', '#80A080', '둥글고 작은 슬라임/구체형'),
    ('N02', '척후 마물', 'Quick Node', '#A0A060', '날씬하고 다리가 긴 곤충형'),
    ('N03', '중갑 마물', 'Heavy Node', '#705050', '크고 넓은 골렘/거인형'),
    ('N04', '갑주 마물', 'Shield Node', '#6080A0', '방패를 든 기사형'),
    ('N05', '벌레 무리', 'Swarm Node', '#908070', '극소형, 벌레/쥐 떼'),
    ('N06', '재생 마물', 'Regen Node', '#508050', '식물/이끼 덮인 유기체'),
    ('N07', '유령 마물', 'Phase Node', '#706090', '반투명, 유령형'),
    ('N08', '분열 마물', 'Split Node', '#806050', '불안정한 점토/분열체'),
    ('N09', '대마 (보스)', 'Boss Node', '#A04040', '다른 몬스터 2~3배 크기, 발광하는 눈'),
]

for i, (mid, name, eng, color, desc) in enumerate(monsters):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.1)
    y = Inches(1.1 + row * 2.0)

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(3.8), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x28)
    card.line.fill.solid()
    card.line.fill.fore_color.rgb = FRAME_COLOR

    # ID & name
    add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.5), Inches(0.35),
                f'{mid}  {name}', 14, TEXT_ACCENT, True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.45), Inches(2.5), Inches(0.25),
                eng, 11, TEXT_SUB, False)

    # Color swatch
    add_color_rect(slide, x + Inches(2.8), y + Inches(0.15), Inches(0.8), Inches(0.6), color)
    add_textbox(slide, x + Inches(2.7), y + Inches(0.8), Inches(1.0), Inches(0.25),
                color, 9, TEXT_SUB, False, PP_ALIGN.CENTER)

    # Description
    add_textbox(slide, x + Inches(0.15), y + Inches(1.1), Inches(3.5), Inches(0.6),
                desc, 11, TEXT_MAIN, False)


# ============================================================
# SLIDE 8: 맵/배경 + 스테이지별 색조
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '07  MAP / BACKGROUND', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

# Layer structure
add_textbox(slide, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
            '레이어 구조 (뒤 → 앞)', 16, TEXT_ACCENT, True)

layers = [
    ('L0', '단색 배경', '스테이지별 색조'),
    ('L1', '미니멀 패턴 타일맵', '경로/배치칸 구분'),
    ('L2', '장식 오브젝트', '스테이지 분위기 전달'),
    ('L3', '경로 하이라이트', '몬스터 이동 경로'),
    ('L4', '타워/몬스터/이펙트', '주인공 레이어'),
    ('L5', 'UI 오버레이', '게임 정보'),
]

for i, (lid, name, desc) in enumerate(layers):
    y = Inches(1.6 + i * 0.38)
    # Gradient from dark to light
    brightness = 0x10 + i * 0x0A
    c = RGBColor(brightness, brightness, brightness + 0x10)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), y, Inches(0.5), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()

    add_textbox(slide, Inches(1.4), y, Inches(1.5), Inches(0.3),
                f'{lid}: {name}', 11, TEXT_MAIN, True)
    add_textbox(slide, Inches(3.2), y, Inches(3), Inches(0.3),
                desc, 11, TEXT_SUB, False)

# Stage colors
add_textbox(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
            '스테이지별 환경 색조 (10개 스테이지)', 16, TEXT_ACCENT, True)

stages = [
    ('1', '잊혀진 초원', '#102020'),
    ('2', '무너진 성채', '#181828'),
    ('3', '수정 동굴', '#1A1030'),
    ('4', '갈라진 협곡', '#201810'),
    ('5', '마탑 내부', '#201020'),
    ('6', '영혼의 숲', '#101A10'),
    ('7', '심연의 사원', '#140A20'),
    ('8', '화산 지대', '#200808'),
    ('9', '왕좌의 방', '#1A1408'),
    ('10', '차원의 균열', '#1A1020'),
]

for i, (num, name, hex_c) in enumerate(stages):
    x = Inches(0.5 + i * 1.25)
    add_color_rect(slide, x, Inches(4.5), Inches(1.05), Inches(1.0), hex_c)
    add_multiline_text(slide, x, Inches(5.55), Inches(1.05), Inches(0.9), [
        (f'S{num}', 10, TEXT_ACCENT, True),
        (name, 9, TEXT_MAIN, False),
        (hex_c, 8, TEXT_SUB, False),
    ])


# ============================================================
# SLIDE 9: UI 스타일
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '08  UI STYLE', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

# UI principles
ui_principles = [
    ('정보 우선', '모든 수치, 상태 정보는 즉시 읽힘'),
    ('판타지 프레임', '패널 테두리에 석재/금속/룬 문양 장식'),
    ('낡은 양피지 톤', '텍스트 배경에 약간의 텍스처 느낌'),
    ('최소 컬러', '무채색+금색 기본, 기능별 색상은 아이콘/강조에만'),
]

for i, (title, desc) in enumerate(ui_principles):
    y = Inches(1.2 + i * 0.5)
    add_textbox(slide, Inches(0.8), y, Inches(2.5), Inches(0.4),
                title, 14, TEXT_ACCENT, True)
    add_textbox(slide, Inches(3.5), y, Inches(8.5), Inches(0.4),
                desc, 14, TEXT_MAIN, False)

# UI Palette
add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12), Inches(0.4),
            'UI 팔레트', 16, TEXT_ACCENT, True)

ui_colors = [
    ('#12101A', 'UI 배경'),
    ('#1A1828', 'UI 패널'),
    ('#242236', '밝은 패널'),
    ('#5A5070', '프레임 기본'),
    ('#B0A080', '프레임 활성'),
    ('#FFD84D', '프레임 강조'),
    ('#E0DCD0', '텍스트 메인'),
    ('#A09890', '텍스트 서브'),
    ('#605850', '텍스트 비활성'),
]

for i, (hex_c, label) in enumerate(ui_colors):
    x = Inches(0.5 + i * 1.38)
    add_color_rect(slide, x, Inches(3.9), Inches(1.2), Inches(0.7), hex_c)
    add_multiline_text(slide, x, Inches(4.65), Inches(1.2), Inches(0.6), [
        (label, 10, TEXT_MAIN, False),
        (hex_c, 9, TEXT_SUB, False),
    ])

# UI Asset list (compact)
add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.4),
            'UI 에셋 목록', 16, TEXT_ACCENT, True)

ui_assets = [
    '패널 프레임 (9-slice, 64x64)',
    '버튼 프레임 기본/강조 (4상태)',
    'HP 바 프레임+필 (200x16)',
    '스킬 노드 (3상태, 48x48)',
    '스킬 연결선 (2상태)',
    '아이콘: Bit(마나), Core(룬석)',
    '타워 인벤토리 슬롯 (48x48)',
    '툴팁/토스트 프레임',
]

for i, asset in enumerate(ui_assets):
    col = i % 2
    row = i // 2
    add_textbox(slide, Inches(0.8 + col * 6), Inches(5.9 + row * 0.35),
                Inches(5.8), Inches(0.35), f'• {asset}', 11, TEXT_MAIN, False)


# ============================================================
# SLIDE 10: VFX / 이펙트
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '09  VFX / EFFECTS', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

add_textbox(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
            '마법 이펙트 = 시각적 주인공  |  타워 시그니처 컬러 준수  |  성장할수록 스케일업',
            14, TEXT_MAIN, False)

# 5-tier effect density
tiers = [
    ('마이크로', '매초', '#4080D4',
     '몬스터 처치 파편, Bit 드롭 팝업,\n투사체 타격 스파크, 투사체 궤적'),
    ('소형', '30초~1분', '#40D470',
     '웨이브 클리어 플래시, WAVE CLEAR 배너,\n보너스 Bit 연출'),
    ('중형', '런 사이', '#FFD84D',
     '스킬 노드 해금, 수치 비교 팝업,\n신규 노드 등장'),
    ('대형', '30분~1시간', '#E08030',
     '스테이지 클리어, Core 획득,\n신규 스테이지 해금'),
    ('특대', '1시간+', '#D44040',
     '보스 처치 (슬로우모션+폭발),\n신규 타워 첫 배치, 엔딩 연출'),
]

for i, (tier, freq, color, desc) in enumerate(tiers):
    x = Inches(0.3 + i * 2.55)
    y = Inches(1.7)

    # Tier card
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(2.35), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x28)
    card.line.fill.solid()
    card.line.fill.fore_color.rgb = hex_to_rgb(color)
    card.line.width = Pt(2)

    # Color bar on top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x, y, Inches(2.35), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(color)
    bar.line.fill.background()

    add_textbox(slide, x + Inches(0.1), y + Inches(0.2), Inches(2.15), Inches(0.35),
                tier, 18, hex_to_rgb(color), True, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.6), Inches(2.15), Inches(0.25),
                f'빈도: {freq}', 11, TEXT_SUB, False, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(1.0), Inches(2.15), Inches(2.3),
                desc, 11, TEXT_MAIN, False)

# Tech spec at bottom
add_multiline_text(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.8), [
    ('기술 사양', 14, TEXT_ACCENT, True),
    ('파티클: Unity Particle System (URP 2D)  |  스프라이트 이펙트: 4~8프레임, 12fps  |  풀링: 최대 동시 50개', 11, TEXT_MAIN, False),
    ('숫자 팝업: TMP + DOTween  |  화면 이펙트: URP 2D Renderer Feature + PostProcessing (Bloom, Vignette)', 11, TEXT_MAIN, False),
])


# ============================================================
# SLIDE 11: 픽셀 아트 스펙
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '10  PIXEL ART SPEC', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

# Resolution table
add_textbox(slide, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
            '스프라이트 해상도', 16, TEXT_ACCENT, True)

res_items = [
    ('타워 (일반)', '64x64 px', 'PPU 64'),
    ('타워 (강화)', '80x80 px', 'PPU 64'),
    ('몬스터 (소형)', '32x32 px', 'PPU 64'),
    ('몬스터 (중형)', '48x48 px', 'PPU 64'),
    ('몬스터 (보스)', '64x64 px', 'PPU 64'),
    ('Swarm', '16x16 px', 'PPU 64'),
    ('타일맵', '32x32 px', 'PPU 32'),
    ('UI 아이콘', '24~48 px', 'Canvas 기준'),
]

for i, (cat, size, ppu) in enumerate(res_items):
    y = Inches(1.6 + i * 0.35)
    add_textbox(slide, Inches(0.8), y, Inches(2.5), Inches(0.3),
                cat, 12, TEXT_MAIN, False)
    add_textbox(slide, Inches(3.5), y, Inches(1.5), Inches(0.3),
                size, 12, TEXT_ACCENT, True)
    add_textbox(slide, Inches(5.2), y, Inches(1.5), Inches(0.3),
                ppu, 12, TEXT_SUB, False)

# Palette limits
add_textbox(slide, Inches(7), Inches(1.15), Inches(5), Inches(0.4),
            '팔레트 제한', 16, TEXT_ACCENT, True)

palette_limits = [
    ('전체 마스터 팔레트', '최대 48색'),
    ('개별 스프라이트', '8~16색'),
    ('환경 타일', '4~6색/스테이지'),
    ('이펙트', '4~8색'),
]

for i, (item, limit) in enumerate(palette_limits):
    y = Inches(1.6 + i * 0.4)
    add_textbox(slide, Inches(7.2), y, Inches(2.5), Inches(0.35),
                item, 12, TEXT_MAIN, False)
    add_textbox(slide, Inches(10), y, Inches(2), Inches(0.35),
                limit, 12, TEXT_ACCENT, True)

# Animation spec
add_textbox(slide, Inches(7), Inches(3.3), Inches(5), Inches(0.4),
            '애니메이션 규격', 16, TEXT_ACCENT, True)

anim_items = [
    ('기본 프레임 레이트', '8fps'),
    ('빠른 액션', '12fps'),
    ('이펙트', '12fps'),
    ('최소 루프', '2프레임'),
    ('최대 시퀀스', '8프레임'),
]

for i, (item, val) in enumerate(anim_items):
    y = Inches(3.8 + i * 0.35)
    add_textbox(slide, Inches(7.2), y, Inches(2.5), Inches(0.3),
                item, 12, TEXT_MAIN, False)
    add_textbox(slide, Inches(10), y, Inches(2), Inches(0.3),
                val, 12, TEXT_ACCENT, True)

# Light rules
add_textbox(slide, Inches(0.5), Inches(4.6), Inches(6), Inches(0.4),
            '광원 규칙', 16, TEXT_ACCENT, True)

light_items = [
    '주 광원 방향: 좌상단 (45도)',
    '하이라이트: 스프라이트 좌상단',
    '그림자: 우하단',
    '환경광: 거의 없음 (다크 판타지)',
    '마법 발광: 자체 발광 (전 방향)',
]

for i, item in enumerate(light_items):
    add_textbox(slide, Inches(0.8), Inches(5.0 + i * 0.35), Inches(5.5), Inches(0.3),
                f'• {item}', 12, TEXT_MAIN, False)

# Performance budget
add_textbox(slide, Inches(7), Inches(5.6), Inches(5), Inches(0.4),
            '성능 예산', 16, TEXT_ACCENT, True)

perf_items = [
    ('동시 스프라이트', '200개 이하'),
    ('드로우콜', '50 이하'),
    ('파티클', '동시 500개 이하'),
    ('2D Light', '동시 25개 이하'),
    ('타겟 FPS', '60fps (최소 30)'),
]

for i, (item, val) in enumerate(perf_items):
    y = Inches(6.0 + i * 0.28)
    add_textbox(slide, Inches(7.2), y, Inches(2.5), Inches(0.25),
                item, 11, TEXT_MAIN, False)
    add_textbox(slide, Inches(10), y, Inches(2), Inches(0.25),
                val, 11, TEXT_ACCENT, True)


# ============================================================
# SLIDE 12: 제작 로드맵
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '11  PRODUCTION ROADMAP', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

phases = [
    ('Phase 1', '코어 루프 (최우선)', '#D44040', [
        'Arrow Tower (Lv1)',
        'Bit/Quick/Heavy Node (이동+파괴)',
        '기본 타일셋 (스테이지 1)',
        '투사체 (Arrow) + 타격',
        'Bit 아이콘 + 숫자 팝업',
        'HP 바 프레임 + 필',
    ]),
    ('Phase 2', 'Hub + 스킬 트리', '#E08030', [
        '스킬 노드 아이콘 (18종)',
        '스킬 트리 프레임/연결선',
        'UI 패널/버튼 프레임',
        'Core 아이콘',
        '타이틀 로고 (판타지)',
    ]),
    ('Phase 3', '추가 타워+몬스터', '#FFD84D', [
        'Cannon/Ice Tower (Lv1)',
        'Shield/Swarm Node',
        '스테이지 2~4 타일셋',
        '타워별 투사체/이펙트',
        '웨이브/스테이지 클리어 연출',
    ]),
    ('Phase 4', '후반 컨텐츠', '#4080D4', [
        'Lightning/Laser/Void Tower',
        'Regen/Phase/Split/Boss Node',
        '스테이지 5~10 타일셋',
        '타워 레벨업 변형 (Lv3, Lv5)',
        '보스 등장/처치 연출',
    ]),
    ('Phase 5', '폴리싱', '#9060D0', [
        '2D Light 세팅 (스테이지별)',
        'Bloom/PostFX 조정',
        '장식 오브젝트 (스테이지별)',
        '타워 아이들 애니메이션',
        '콤보/특수 이펙트',
    ]),
]

for i, (phase, subtitle, color, items) in enumerate(phases):
    x = Inches(0.3 + i * 2.55)
    y = Inches(1.2)

    # Phase card
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(2.35), Inches(5.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x28)
    card.line.fill.solid()
    card.line.fill.fore_color.rgb = hex_to_rgb(color)
    card.line.width = Pt(2)

    # Color header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x, y, Inches(2.35), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(color)
    bar.line.fill.background()

    add_textbox(slide, x + Inches(0.05), y + Inches(0.05), Inches(2.25), Inches(0.4),
                phase, 18, RGBColor(0xFF, 0xFF, 0xFF), True, PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.1), y + Inches(0.6), Inches(2.15), Inches(0.35),
                subtitle, 12, hex_to_rgb(color), True, PP_ALIGN.CENTER)

    items_text = '\n'.join([f'• {item}' for item in items])
    add_textbox(slide, x + Inches(0.1), y + Inches(1.1), Inches(2.15), Inches(4.5),
                items_text, 11, TEXT_MAIN, False)


# ============================================================
# SLIDE 13: 용어 전환표
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '12  TERMINOLOGY TRANSITION', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

add_textbox(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
            '코드명은 유지, 인게임 표시 텍스트만 변경 검토', 14, TEXT_SUB, False)

# Table header
for col_i, (header, x_pos, w) in enumerate([
    ('코드명 (유지)', 0.8, 2.0),
    ('기존 표시', 3.0, 2.5),
    ('신규 표시 후보', 5.7, 3.5),
    ('총괄PD 확인', 9.5, 2.5),
]):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x_pos), Inches(1.7), Inches(w), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x30, 0x20, 0x40)
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = FRAME_COLOR
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_ACCENT
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER

terms = [
    ('Node', 'Node', '마물 / 이형체 / Node(유지)', '예'),
    ('Bit', 'Bit', '소울 / 마나 / Bit(유지)', '예'),
    ('Core', 'Core', '룬석 / 정수 / Core(유지)', '예'),
    ('Tower', 'Tower', '마탑 / Tower(유지)', '예'),
    ('Hub', 'Hub', '마법진 / Hub(유지)', '예'),
    ('스킬 트리', '스킬 트리', '룬 각성 트리 / 스킬 트리(유지)', '예'),
    ('Wave', '웨이브', '침공 / 웨이브(유지)', '예'),
]

for i, (code, old, new, confirm) in enumerate(terms):
    y = Inches(2.25 + i * 0.45)
    bg_c = RGBColor(0x1A, 0x18, 0x28) if i % 2 == 0 else RGBColor(0x16, 0x14, 0x22)

    for x_pos, w, text, color in [
        (0.8, 2.0, code, TEXT_MAIN),
        (3.0, 2.5, old, TEXT_SUB),
        (5.7, 3.5, new, TEXT_MAIN),
        (9.5, 2.5, confirm, TEXT_ACCENT),
    ]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x_pos), y, Inches(w), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_c
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = RGBColor(0x30, 0x28, 0x40)
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER

# Game name note
add_multiline_text(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(1.2), [
    ('게임명 확정', 16, TEXT_ACCENT, True),
    ('', 6),
    ('기존 "Nodebreaker TD" → 신규 "Soulspire"로 확정', 14, TEXT_MAIN, True),
    ('세계관: 다크 판타지  |  장르: 픽셀 아트 타워 디펜스', 13, TEXT_SUB, False),
])


# ============================================================
# SLIDE 14: 오픈 이슈
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            '13  OPEN ISSUES', 28, TEXT_ACCENT, True)
add_divider(slide, Inches(0.95))

issues = [
    ('1', 'AI 이미지 생성 도구 최종 선정', 'TA팀장', '미정'),
    ('2', 'AI 사운드 생성 도구 최종 선정', 'TA팀장', '미정'),
    ('3', '다크 판타지 스타일 락 실행 (타워)', 'TA팀장', '대기'),
    ('4', '다크 판타지 스타일 락 실행 (몬스터)', 'TA팀장', '대기'),
    ('5', '다크 판타지 스타일 락 실행 (타일셋)', 'TA팀장', '대기'),
    ('6', 'UIColors.cs 팔레트 업데이트', '프로그래밍팀장', '기획안 확정 후'),
    ('7', '기존 Art 폴더 아카이브 처리', 'TA팀장', '기획안 확정 후'),
    ('8', '타일맵 셀 크기 최종 검증', 'PG+TA팀장', '대기'),
    ('9', '인게임 텍스트 변경 범위', '기획팀장', '총괄PD 결정 필요'),
    ('10', '게임명 확정 (→ Soulspire)', '총괄PD', '확정'),
]

# Table headers
for x_pos, w, header in [
    (0.5, 0.5, '#'),
    (1.1, 5.5, '이슈'),
    (6.8, 2.5, '담당'),
    (9.5, 2.8, '상태'),
]:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x_pos), Inches(1.2), Inches(w), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x30, 0x20, 0x40)
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = FRAME_COLOR
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_ACCENT
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER

for i, (num, issue, owner, status) in enumerate(issues):
    y = Inches(1.75 + i * 0.45)
    bg_c = RGBColor(0x1A, 0x18, 0x28) if i % 2 == 0 else RGBColor(0x16, 0x14, 0x22)

    status_color = TEXT_ACCENT
    if status == '미정':
        status_color = RGBColor(0xD4, 0x40, 0x40)
    elif status == '대기':
        status_color = RGBColor(0xE0, 0xD0, 0x40)
    elif status == '확정':
        status_color = RGBColor(0x40, 0xD4, 0x70)

    for x_pos, w, text, color, align in [
        (0.5, 0.5, num, TEXT_SUB, PP_ALIGN.CENTER),
        (1.1, 5.5, issue, TEXT_MAIN, PP_ALIGN.LEFT),
        (6.8, 2.5, owner, TEXT_SUB, PP_ALIGN.CENTER),
        (9.5, 2.8, status, status_color, PP_ALIGN.CENTER),
    ]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x_pos), y, Inches(w), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_c
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = RGBColor(0x30, 0x28, 0x40)
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.alignment = align


# ============================================================
# SAVE
# ============================================================
output_path = '/mnt/c/UnityProjects/test1/Docs/Design/Soulspire_ArtDirection_v0.1.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
